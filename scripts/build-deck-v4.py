"""Deck v4 — built from scratch. Clean flow, consistent layout, all screenshots inline."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os
from PIL import Image as PILImage

OUT = "/Users/saikanagat/Desktop/backpack-bottle/Backpack-Bottle-Campaign-v4.pptx"
SHOTS = "/Users/saikanagat/Desktop/backpack-bottle/assets/deck"

# Brand palette
FOREST = RGBColor(0x1F, 0x3A, 0x2E)
FOREST_DARK = RGBColor(0x14, 0x26, 0x20)
CREAM = RGBColor(0xF5, 0xEF, 0xE6)
CREAM_DARK = RGBColor(0xE8, 0xE0, 0xD0)
BURNT = RGBColor(0xD9, 0x76, 0x42)
BURNT_DARK = RGBColor(0xB8, 0x5F, 0x31)
INK = RGBColor(0x1F, 0x3A, 0x2E)
MUTED = RGBColor(0x5A, 0x6B, 0x62)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Build new deck — 16:9 default
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]
PAGE_TOTAL_PLACEHOLDER = "{TOTAL}"

# ---------- helpers ----------------------------------------------
def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False,
                color=INK, font_name="Helvetica", align=PP_ALIGN.LEFT, italic=False,
                line_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing: p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb

def add_multiline(slide, left, top, width, height, lines, font_size=12,
                  color=INK, font_name="Helvetica", align=PP_ALIGN.LEFT, line_spacing=1.3):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        # support tuple (text, bold, color, size_override) for inline emphasis
        if isinstance(line, tuple):
            text, bold = line[0], line[1] if len(line) > 1 else False
            c = line[2] if len(line) > 2 else color
            sz = line[3] if len(line) > 3 else font_size
        else:
            text, bold, c, sz = line, False, color, font_size
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = c
        run.font.name = font_name
    return tb

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width: shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_filled_bg(slide, color):
    bg = add_rect(slide, 0, 0, SW, SH, color)
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg

def add_eyebrow(slide, color=BURNT):
    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.3),
                "GROUP 3 · DIGITAL PLATFORMS LAB · BOLOGNA BUSINESS SCHOOL",
                font_size=8.5, bold=True, color=color)

def add_footer(slide, page, total_placeholder=True):
    add_textbox(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
                "backpack & bottle  ·  Group 3", font_size=9, color=MUTED)
    total_str = PAGE_TOTAL_PLACEHOLDER if total_placeholder else ""
    add_textbox(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
                f"{page}  /  {total_str}", font_size=9, color=MUTED, align=PP_ALIGN.RIGHT)

def add_title(slide, eyebrow_text, title, sub=None, top=Inches(0.95)):
    add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.35),
                eyebrow_text.upper(), font_size=10, bold=True, color=BURNT)
    add_textbox(slide, Inches(0.6), top + Inches(0.4), Inches(12), Inches(0.7),
                title, font_size=30, bold=False, color=FOREST,
                font_name="Times-Bold")
    if sub:
        add_textbox(slide, Inches(0.6), top + Inches(1.1), Inches(12), Inches(0.4),
                    sub, font_size=12, color=MUTED, italic=True)

def new_content_slide(eyebrow, title, sub=None):
    s = prs.slides.add_slide(BLANK)
    add_filled_bg(s, CREAM)
    add_eyebrow(s)
    add_title(s, eyebrow, title, sub)
    return s

def new_section_divider(num, label, sub):
    """Full-bleed forest divider with big number + label."""
    s = prs.slides.add_slide(BLANK)
    add_filled_bg(s, FOREST)
    # Big number
    add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(2),
                num, font_size=130, bold=True, color=BURNT,
                font_name="Times-Bold")
    add_textbox(s, Inches(1), Inches(4.6), Inches(11), Inches(0.8),
                label, font_size=42, color=CREAM, font_name="Times-Bold")
    add_textbox(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                sub, font_size=14, color=BURNT, italic=True)
    add_textbox(s, Inches(1), Inches(0.5), Inches(11), Inches(0.3),
                "GROUP 3 · DIGITAL PLATFORMS LAB · BOLOGNA BUSINESS SCHOOL",
                font_size=8.5, bold=True, color=BURNT)
    return s

def add_image_below_title(slide, image_path, top_in=2.4, max_w_in=11.5, max_h_in=4.4,
                          centered=True, with_border=True):
    if not os.path.exists(image_path): return None
    img = PILImage.open(image_path)
    iw, ih = img.size
    aspect = ih / iw
    final_w = min(max_w_in, max_h_in / aspect)
    final_h = final_w * aspect
    if final_h > max_h_in:
        final_h = max_h_in
        final_w = final_h / aspect
    left = (13.333 - final_w) / 2 if centered else 0.6
    if with_border:
        # Drop a thin burnt border
        add_rect(slide, Inches(left - 0.04), Inches(top_in - 0.04),
                 Inches(final_w + 0.08), Inches(final_h + 0.08),
                 CREAM_DARK)
    pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top_in),
                                   width=Inches(final_w), height=Inches(final_h))
    return pic

# Set page totals at end
PAGES_INFO = []  # list of (slide, page_num)
def page_register(slide, num):
    PAGES_INFO.append((slide, num))

# =================================================================
# DECK BUILD
# =================================================================
PAGE = 0

# ---------- 1. COVER --------------------------------------------
def slide_cover():
    global PAGE; PAGE += 1
    s = prs.slides.add_slide(BLANK)
    add_filled_bg(s, FOREST)
    # Try to add a subtle hero image strip on the right
    if os.path.exists(os.path.join(SHOTS, "hero-cta.png")):
        try:
            img = PILImage.open(os.path.join(SHOTS, "hero-cta.png"))
            iw, ih = img.size
            tg_w_in = 5.0
            ar = ih / iw
            tg_h_in = tg_w_in * ar
            if tg_h_in > 7.5:
                tg_h_in = 7.5; tg_w_in = tg_h_in / ar
            s.shapes.add_picture(os.path.join(SHOTS, "hero-cta.png"),
                                 Inches(13.333 - tg_w_in), Inches(0),
                                 width=Inches(tg_w_in), height=Inches(tg_h_in))
            # Dark overlay over the image (right 5 inches)
            ov = add_rect(s, Inches(13.333 - tg_w_in), Inches(0),
                          Inches(tg_w_in), Inches(7.5), FOREST_DARK)
            ov.fill.solid()
            ov.fill.fore_color.rgb = FOREST_DARK
            ov.fill.transparency = 0.5  # type: ignore
        except Exception:
            pass

    add_textbox(s, Inches(0.6), Inches(0.5), Inches(8), Inches(0.3),
                "GROUP 3 · DIGITAL PLATFORMS LAB · BOLOGNA BUSINESS SCHOOL",
                font_size=9, bold=True, color=BURNT)
    add_textbox(s, Inches(0.6), Inches(2.1), Inches(11), Inches(1.4),
                "backpack & bottle.",
                font_size=72, color=CREAM, font_name="Times-Bold")
    add_textbox(s, Inches(0.6), Inches(3.5), Inches(11), Inches(0.6),
                "Curated European city breaks. Transparent prices. €50 off.",
                font_size=20, color=CREAM, italic=True)
    add_textbox(s, Inches(0.6), Inches(4.1), Inches(11), Inches(0.5),
                "Campaign brief · brand · landing page · tracking · media plan · automation.",
                font_size=13, color=BURNT)
    # Team
    add_textbox(s, Inches(0.6), Inches(5.5), Inches(2.5), Inches(0.3),
                "GROUP 3 · TEAM", font_size=10, bold=True, color=BURNT)
    add_textbox(s, Inches(0.6), Inches(5.85), Inches(11), Inches(0.4),
                "Sai Prathyaksh Kanagat  ·  Marica Motta  ·  Alberto Faggiotto",
                font_size=12, color=CREAM)
    add_textbox(s, Inches(0.6), Inches(6.2), Inches(11), Inches(0.4),
                "Cecilia Gullett  ·  Marilu Stevens  ·  Amelia Ruben",
                font_size=12, color=CREAM)
    add_textbox(s, Inches(0.6), Inches(6.85), Inches(11), Inches(0.3),
                "Campaign window 5 May – 15 June 2026  ·  Presented Tuesday 5 May 2026",
                font_size=10, color=BURNT, italic=True)
slide_cover()

# ---------- 2. AGENDA --------------------------------------------
def slide_agenda():
    global PAGE; PAGE += 1
    s = new_content_slide("Agenda", "What we'll walk through",
                          "Twelve sections · ~40 minutes · questions held to the end")
    items = [
        ("01", "The brief", "What we were asked to do"),
        ("02", "Strategic decisions", "Audience · offer · positioning"),
        ("03", "Brand & offer", "Identity, 5 cities, BB50 coupon"),
        ("04", "The product", "Live landing page · mobile · destination detail"),
        ("05", "Tracking architecture", "GTM · GA4 · dataLayer schema"),
        ("06", "Compliance", "Consent Mode v2 · cookie tracking · GDPR"),
        ("07", "Email & automation", "Brevo · welcome series · form pipeline"),
        ("08", "Media plan", "Meta · Google · budget · pacing"),
        ("09", "KPIs & measurement", "Funnel targets · Looker Studio"),
        ("10", "Optimization", "A/B test · competitors · accessibility"),
        ("11", "Deliverables", "Vercel · GitHub · live URLs"),
        ("12", "Closing", "Risks · what's next · thank you"),
    ]
    y0 = 2.4
    for i, (n, title, sub) in enumerate(items):
        col = i // 6
        row = i % 6
        x = 0.6 + col * 6.3
        yy = y0 + row * 0.65
        add_textbox(s, Inches(x), Inches(yy), Inches(0.6), Inches(0.4),
                    n, font_size=14, bold=True, color=BURNT, font_name="Times-Bold")
        add_textbox(s, Inches(x + 0.7), Inches(yy + 0.02), Inches(5.4), Inches(0.3),
                    title, font_size=14, bold=True, color=FOREST)
        add_textbox(s, Inches(x + 0.7), Inches(yy + 0.32), Inches(5.4), Inches(0.3),
                    sub, font_size=10, color=MUTED)
    add_footer(s, PAGE)
slide_agenda()

# ===== SECTION 01 ============================================
new_section_divider("01", "The brief", "What we were asked to do · Group 3 · Backpack & Bottle"); PAGE += 1

# ---- THE BRIEF SLIDE
def slide_brief():
    global PAGE; PAGE += 1
    s = new_content_slide("The brief", "Group 3 — Backpack & Bottle",
                          "Promote travel deals. Generate qualified leads. Italian millennials.")
    cards = [
        ("BUSINESS OBJECTIVE",
         "Promote travel deals under the 'Backpack & Bottle' brand and generate qualified leads interested in affordable, experience-driven European city breaks."),
        ("MARKETING OBJECTIVE",
         "Make the offer feel attractive, accessible, and time-sensitive — combining inspiration with a clear incentive to request the BB50 coupon."),
        ("WHAT WAS GIVEN",
         "Brand name. Audience demographic. Offer mechanic outline. €2,000 budget. 6-week window. 5 destinations to pick. 60/40 Meta/Google split."),
        ("WHAT WAS OPEN",
         "Brand identity, voice, positioning. Specific cities. Coupon code/value. Landing page. Tracking infra. Email automation. KPI targets."),
    ]
    y = 2.5
    for i, (label, text) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = 0.6 + col * 6.3
        yy = y + row * 1.95
        add_rect(s, Inches(x), Inches(yy), Inches(6.05), Inches(1.75), WHITE,
                 line_color=CREAM_DARK)
        add_textbox(s, Inches(x + 0.3), Inches(yy + 0.2), Inches(5.5), Inches(0.3),
                    label, font_size=10, bold=True, color=BURNT)
        add_textbox(s, Inches(x + 0.3), Inches(yy + 0.55), Inches(5.5), Inches(1.1),
                    text, font_size=11, color=INK, line_spacing=1.3)
    add_footer(s, PAGE)
slide_brief()

# ---- STRATEGIC DECISIONS GLANCE
def slide_decisions():
    global PAGE; PAGE += 1
    s = new_content_slide("Strategic decisions at a glance", "Six choices that shape every other choice",
                          "Audience · offer · coupon · channels · KPIs · timing")
    decisions = [
        ("AUDIENCE", "Italian millennials 25–39\nUrban · €25–45k income"),
        ("OFFER", "5 curated European cities\n€280–365 · voli + hotel"),
        ("COUPON", "BB50 · €50 off · 90 days\nMin €300 booking"),
        ("CHANNELS", "60% Meta · 40% Google\nAwareness + Lead Gen + Retarget"),
        ("KPIs", "5,000 visits · 300 leads\n€6.67 CPL · 1.8x ROAS proj"),
        ("TIMING", "5 May – 15 June 2026\n6 weeks · pre-summer ramp"),
    ]
    y = 2.5
    for i, (label, text) in enumerate(decisions):
        col = i % 3
        row = i // 3
        x = 0.6 + col * 4.18
        yy = y + row * 1.95
        add_rect(s, Inches(x), Inches(yy), Inches(3.95), Inches(1.75),
                 FOREST if i % 2 == 0 else CREAM_DARK)
        c_label = BURNT if i % 2 == 0 else BURNT
        c_text = CREAM if i % 2 == 0 else INK
        add_textbox(s, Inches(x + 0.25), Inches(yy + 0.2), Inches(3.5), Inches(0.3),
                    label, font_size=11, bold=True, color=c_label)
        add_multiline(s, Inches(x + 0.25), Inches(yy + 0.6), Inches(3.5), Inches(1.1),
                      text.split("\n"), font_size=12, color=c_text, line_spacing=1.4)
    add_footer(s, PAGE)
slide_decisions()

# ---- TEAM & ROLES
def slide_team():
    global PAGE; PAGE += 1
    s = new_content_slide("Group 3 · Team", "Six roles, one shipped campaign",
                          "Each speaker owns the section their role built")
    roles = [
        ("01", "Sai Prathyaksh Kanagat", "CRO Specialist + Lead",
         "Landing page, brand strategy, conversion optimization, project lead."),
        ("02", "Marica Motta", "MarTech Specialist",
         "GTM container, GA4 events, Meta Pixel, dataLayer."),
        ("03", "Alberto Faggiotto", "Compliance Specialist",
         "Cookie banner, Consent Mode v2, GDPR privacy, T&Cs."),
        ("04", "Cecilia Gullett", "Email Marketing",
         "Brevo setup, 8-email welcome series, contact attributes."),
        ("05", "Marilu Stevens", "Marketing Automation",
         "Form → Brevo bridge, Meta CAPI, UTM persistence."),
        ("06", "Amelia Ruben", "Campaign Analyst",
         "KPI tree, custom dimensions, Looker Studio dashboard."),
    ]
    y = 2.45
    for i, (num, name, role, scope) in enumerate(roles):
        col = i % 2
        row = i // 2
        x = 0.6 + col * 6.18
        yy = y + row * 1.45
        add_rect(s, Inches(x), Inches(yy), Inches(5.95), Inches(1.3), WHITE,
                 line_color=CREAM_DARK)
        add_rect(s, Inches(x), Inches(yy), Inches(0.6), Inches(1.3), BURNT)
        add_textbox(s, Inches(x + 0.05), Inches(yy + 0.5), Inches(0.5), Inches(0.4),
                    num, font_size=18, bold=True, color=CREAM, font_name="Times-Bold",
                    align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(x + 0.8), Inches(yy + 0.15), Inches(5.0), Inches(0.4),
                    name, font_size=14, bold=True, color=FOREST, font_name="Times-Bold")
        add_textbox(s, Inches(x + 0.8), Inches(yy + 0.5), Inches(5.0), Inches(0.3),
                    role, font_size=10, bold=True, color=BURNT_DARK)
        add_textbox(s, Inches(x + 0.8), Inches(yy + 0.8), Inches(5.0), Inches(0.5),
                    scope, font_size=10, color=INK)
    add_footer(s, PAGE)
slide_team()

# ===== SECTION 02 ===========================================
new_section_divider("02", "Strategy", "Audience · pain points · positioning · brand"); PAGE += 1

# ---- AUDIENCE
def slide_audience():
    global PAGE; PAGE += 1
    s = new_content_slide("Audience", "Italian millennials 25–39",
                          "Urban · digitally native · time-poor · budget-conscious but not cheap")
    panels = [
        ("DEMOGRAPHICS",
         ["Age 25–39", "Milan · Rome · Bologna · Florence · Turin",
          "Income €25–45k / year", "University degree", "Single or DINK couple"]),
        ("DIGITAL HABITS",
         ["3 hours/day on Instagram + Facebook", "Discovers brands on TikTok",
          "Books mostly on mobile", "Compares 5+ sites before booking",
          "Reads reviews on Booking + TripAdvisor"]),
        ("BOOKING BEHAVIOUR",
         ["Travels 3–5 times/year", "Average trip €250–500",
          "Books 2–8 weeks ahead", "Prefers experiences over luxury",
          "Cares about authenticity, not 5★ resorts"]),
    ]
    y = 2.5
    for i, (label, items) in enumerate(panels):
        x = 0.6 + i * 4.18
        add_rect(s, Inches(x), Inches(y), Inches(3.95), Inches(4.2), WHITE,
                 line_color=CREAM_DARK)
        add_rect(s, Inches(x), Inches(y), Inches(3.95), Inches(0.55), FOREST)
        add_textbox(s, Inches(x + 0.25), Inches(y + 0.18), Inches(3.5), Inches(0.3),
                    label, font_size=10, bold=True, color=BURNT)
        for j, it in enumerate(items):
            add_textbox(s, Inches(x + 0.3), Inches(y + 0.85 + j*0.55), Inches(0.3), Inches(0.3),
                        "·", font_size=18, bold=True, color=BURNT)
            add_textbox(s, Inches(x + 0.55), Inches(y + 0.9 + j*0.55), Inches(3.4), Inches(0.5),
                        it, font_size=11, color=INK, line_spacing=1.2)
    add_footer(s, PAGE)
slide_audience()

# ---- PAIN POINTS
def slide_pain():
    global PAGE; PAGE += 1
    s = new_content_slide("Pain points", "What blocks them from booking",
                          "Five frictions the campaign solves with curation, transparency, and a real incentive")
    pains = [
        ("01", "Research overwhelm",
         "Booking.com shows 500 hotels for one weekend. Paralysing. We show 3–5 curated options per city — pre-vetted for location, reviews, value."),
        ("02", "Price transparency",
         "Hidden fees discovered at checkout. We show all-in prices upfront: flight + hotel + tasting itinerary, no surprise add-ons."),
        ("03", "Trust gap",
         "Generic OTAs feel transactional. We're an editorial brand with a point of view — every recommendation has a story."),
        ("04", "FOMO loop",
         "Friends are travelling, social media is full of weekends. Our coupon (€50) lowers the activation barrier from 'maybe' to 'why not'."),
        ("05", "Time hostage",
         "Comparing 5 sites takes 2 hours. We compress decision-making to 5 minutes — single page, single offer, single click."),
    ]
    y = 2.45
    for i, (num, title, desc) in enumerate(pains):
        col = i % 3
        row = i // 3
        x = 0.6 + col * 4.18
        yy = y + row * 2.1
        add_rect(s, Inches(x), Inches(yy), Inches(3.95), Inches(1.95), WHITE,
                 line_color=CREAM_DARK)
        add_textbox(s, Inches(x + 0.25), Inches(yy + 0.2), Inches(0.6), Inches(0.5),
                    num, font_size=22, bold=True, color=BURNT, font_name="Times-Bold")
        add_textbox(s, Inches(x + 0.95), Inches(yy + 0.27), Inches(2.9), Inches(0.4),
                    title, font_size=14, bold=True, color=FOREST, font_name="Times-Bold")
        add_textbox(s, Inches(x + 0.25), Inches(yy + 0.85), Inches(3.5), Inches(1.05),
                    desc, font_size=10, color=INK, line_spacing=1.35)
    add_footer(s, PAGE)
slide_pain()

# ---- POSITIONING
def slide_positioning():
    global PAGE; PAGE += 1
    s = new_content_slide("Positioning", "How we differ from every competitor in their head",
                          "One sentence. Says what we are, what we promise, and who it's for.")
    add_rect(s, Inches(1.5), Inches(2.8), Inches(10.3), Inches(2.5), FOREST)
    add_textbox(s, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.4),
                "OUR POSITIONING", font_size=11, bold=True, color=BURNT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.5), Inches(3.45), Inches(10.3), Inches(1.0),
                '"Skip the research fatigue."',
                font_size=44, color=CREAM, font_name="Times-Bold",
                italic=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.5), Inches(4.65), Inches(10.3), Inches(0.6),
                "European city breaks at insider prices — early-access deals that fit your budget",
                font_size=14, color=BURNT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.4),
                "and deliver experiences worth sharing.",
                font_size=14, color=BURNT, align=PP_ALIGN.CENTER)
    # 3 supports
    supports = [
        ("CURATED", "3–5 picks per city, not 500"),
        ("TRANSPARENT", "All-in price, no hidden fees"),
        ("AUTHENTIC", "Editorial voice, real recommendations"),
    ]
    y_sup = 5.7
    for i, (lbl, txt) in enumerate(supports):
        x = 1.5 + i * 3.45
        add_textbox(s, Inches(x), Inches(y_sup), Inches(3.4), Inches(0.3),
                    lbl, font_size=10, bold=True, color=BURNT, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(x), Inches(y_sup + 0.3), Inches(3.4), Inches(0.4),
                    txt, font_size=12, color=FOREST, align=PP_ALIGN.CENTER)
    add_footer(s, PAGE)
slide_positioning()

# ---- BRAND IDENTITY
def slide_brand():
    global PAGE; PAGE += 1
    s = new_content_slide("Brand identity", "Modern Wanderer — warm, editorial, trustworthy",
                          "Palette · typography · voice · all wired into the codebase as design tokens")
    # Palette
    add_textbox(s, Inches(0.6), Inches(2.4), Inches(2), Inches(0.3),
                "PALETTE", font_size=10, bold=True, color=BURNT)
    swatches = [
        ("Forest", "#1F3A2E", FOREST),
        ("Cream", "#F5EFE6", CREAM),
        ("Burnt", "#D97642", BURNT),
        ("Forest dark", "#142620", FOREST_DARK),
        ("Cream dark", "#E8E0D0", CREAM_DARK),
        ("Burnt dark", "#B85F31", BURNT_DARK),
    ]
    for i, (name, hexv, col) in enumerate(swatches):
        x = 0.6 + i * 1.05
        add_rect(s, Inches(x), Inches(2.8), Inches(0.95), Inches(1.3), col,
                 line_color=CREAM_DARK)
        add_textbox(s, Inches(x), Inches(4.2), Inches(0.95), Inches(0.3),
                    name, font_size=9, bold=True, color=FOREST, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(x), Inches(4.45), Inches(0.95), Inches(0.3),
                    hexv, font_size=8, color=MUTED, align=PP_ALIGN.CENTER, font_name="Courier")

    # Typography
    add_textbox(s, Inches(7.5), Inches(2.4), Inches(2), Inches(0.3),
                "TYPOGRAPHY", font_size=10, bold=True, color=BURNT)
    add_textbox(s, Inches(7.5), Inches(2.75), Inches(5.5), Inches(0.6),
                "Fraunces — display headlines",
                font_size=22, color=FOREST, font_name="Times-Bold", italic=True)
    add_textbox(s, Inches(7.5), Inches(3.4), Inches(5.5), Inches(0.4),
                "Editorial · serif · warm character",
                font_size=10, color=MUTED)
    add_textbox(s, Inches(7.5), Inches(4.0), Inches(5.5), Inches(0.5),
                "Inter — body & UI",
                font_size=22, bold=True, color=FOREST)
    add_textbox(s, Inches(7.5), Inches(4.55), Inches(5.5), Inches(0.4),
                "Geometric sans · neutral · highly legible at small sizes",
                font_size=10, color=MUTED)

    # Voice
    add_textbox(s, Inches(0.6), Inches(5.0), Inches(2), Inches(0.3),
                "VOICE", font_size=10, bold=True, color=BURNT)
    voice = [
        "Direct. Says what it does, not what it 'enables'.",
        "Editorial, not transactional. 'Skip the research fatigue', not 'Book now'.",
        "Confident in Italian. Doesn't translate from English.",
        "No exclamation marks. No urgency theatre. Substance over hype.",
    ]
    for i, line in enumerate(voice):
        add_textbox(s, Inches(0.85), Inches(5.4 + i*0.35), Inches(0.3), Inches(0.3),
                    "·", font_size=18, bold=True, color=BURNT)
        add_textbox(s, Inches(1.1), Inches(5.45 + i*0.35), Inches(11.5), Inches(0.3),
                    line, font_size=11, color=INK)
    add_footer(s, PAGE)
slide_brand()

# ===== SECTION 03 ===========================================
new_section_divider("03", "The offer", "5 cities · BB50 coupon · the funnel"); PAGE += 1

# ---- THE OFFER (5 cities + screenshot)
def slide_offer():
    global PAGE; PAGE += 1
    s = new_content_slide("The offer", "5 curated European city breaks",
                          "1.5–3 hour flights from Italy · transparent all-in pricing")
    cities = [
        ("Rome", "€280", "2 nights"),
        ("Barcelona", "€320", "3 nights"),
        ("Amsterdam", "€365", "3 nights"),
        ("Lisbon", "€345", "3 nights"),
        ("Prague", "€295", "2 nights"),
    ]
    y = 2.5
    cw = 2.36
    gap = 0.13
    total = 5 * cw + 4 * gap
    x0 = (13.333 - total) / 2
    for i, (name, price, nights) in enumerate(cities):
        cx = x0 + i * (cw + gap)
        add_rect(s, Inches(cx), Inches(y), Inches(cw), Inches(2.0), WHITE,
                 line_color=CREAM_DARK)
        add_rect(s, Inches(cx), Inches(y), Inches(cw), Inches(0.6), FOREST)
        add_textbox(s, Inches(cx), Inches(y + 0.18), Inches(cw), Inches(0.3),
                    name.upper(), font_size=14, bold=True, color=CREAM,
                    align=PP_ALIGN.CENTER, font_name="Times-Bold")
        add_textbox(s, Inches(cx), Inches(y + 0.85), Inches(cw), Inches(0.5),
                    price, font_size=28, bold=True, color=FOREST,
                    font_name="Times-Bold", align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(cx), Inches(y + 1.4), Inches(cw), Inches(0.3),
                    nights, font_size=11, color=BURNT_DARK, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(cx), Inches(y + 1.65), Inches(cw), Inches(0.3),
                    "voli + hotel inclusi", font_size=9, color=MUTED, italic=True,
                    align=PP_ALIGN.CENTER)

    # Live screenshot below
    add_textbox(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.3),
                "LIVE — DESTINATION GRID ON HOMEPAGE",
                font_size=9, bold=True, color=BURNT)
    add_image_below_title(s, os.path.join(SHOTS, "destinations-grid.png"),
                          top_in=5.0, max_w_in=11, max_h_in=2.0)
    add_footer(s, PAGE)
slide_offer()

# ---- COUPON MECHANIC
def slide_coupon():
    global PAGE; PAGE += 1
    s = new_content_slide("Coupon mechanic", "BB50 — built for lead capture",
                          "Not a race-to-the-bottom price cut · enough to activate, not enough to destroy margin")
    add_rect(s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5), FOREST_DARK)
    add_textbox(s, Inches(1.5), Inches(2.75), Inches(10.3), Inches(0.5),
                "THE MECHANIC", font_size=11, bold=True, color=BURNT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.0),
                "€50",
                font_size=80, color=CREAM, font_name="Times-Bold", align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.5),
                "off bookings from €300",
                font_size=18, color=CREAM, italic=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.4),
                "Code: BB50 · Valid 90 days · One per customer",
                font_size=12, color=BURNT, align=PP_ALIGN.CENTER)

    rules = [
        ("MIN BASKET", "€300", "Industry avg in IT travel = €322"),
        ("DISCOUNT %", "16.7%", "Below the 20% margin alarm zone"),
        ("VALIDITY", "90 days", "Long enough to plan; short enough to urge"),
        ("BLACKOUT", "23 Dec – 2 Jan\n10–20 Aug", "High-demand windows excluded"),
    ]
    y = 5.4
    for i, (lbl, val, desc) in enumerate(rules):
        x = 1.5 + i * 2.62
        add_textbox(s, Inches(x), Inches(y), Inches(2.5), Inches(0.3),
                    lbl, font_size=9, bold=True, color=BURNT, align=PP_ALIGN.CENTER)
        add_multiline(s, Inches(x), Inches(y + 0.32), Inches(2.5), Inches(0.6),
                      val.split("\n"), font_size=14, color=FOREST, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(x), Inches(y + 1.0), Inches(2.5), Inches(0.5),
                    desc, font_size=9, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
    add_footer(s, PAGE)
slide_coupon()

# ---- THE FUNNEL
def slide_funnel():
    global PAGE; PAGE += 1
    s = new_content_slide("The funnel", "From paid impression to booking",
                          "Six stages · each one a measured event · each one a target")
    stages = [
        ("IMPRESSIONS", "150–180k", "Meta + Google Ads", FOREST),
        ("CLICKS", "5,000", "CPC ~€0.40", FOREST),
        ("FORM VIEWS", "1,500", "30% scroll-to-form", BURNT),
        ("FORM SUBMITS", "300", "★ Key Event · €50/lead", BURNT_DARK),
        ("COUPON DOWNLOADS", "400", "★ Key Event", BURNT_DARK),
        ("BOOKINGS", "30", "10% lead-to-book · €322 AOV", FOREST_DARK),
    ]
    y = 3.0
    bw = 1.95
    gap = 0.05
    total_w = 6 * bw + 5 * gap
    x0 = (13.333 - total_w) / 2
    for i, (lbl, val, sub, c) in enumerate(stages):
        x = x0 + i * (bw + gap)
        add_rect(s, Inches(x), Inches(y), Inches(bw), Inches(2.0), c)
        add_textbox(s, Inches(x), Inches(y + 0.2), Inches(bw), Inches(0.3),
                    lbl, font_size=8.5, bold=True, color=BURNT, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(x), Inches(y + 0.65), Inches(bw), Inches(0.6),
                    val, font_size=22, color=CREAM, font_name="Times-Bold",
                    align=PP_ALIGN.CENTER, bold=True)
        add_textbox(s, Inches(x), Inches(y + 1.35), Inches(bw), Inches(0.6),
                    sub, font_size=8.5, color=CREAM, align=PP_ALIGN.CENTER, italic=True)
        if i < len(stages) - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     Inches(x + bw - 0.1), Inches(y + 0.85),
                                     Inches(0.25), Inches(0.3))
            ar.fill.solid(); ar.fill.fore_color.rgb = BURNT
            ar.line.fill.background()
    add_textbox(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.4),
                "Conversion rates: Click→FormView 30% · FormView→Submit 20% · Submit→Coupon 133% (incl. /grazie repeats) · Lead→Booking 10%",
                font_size=10, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.4),
                "Result: €2,000 spend · 30 bookings · €322 AOV → €9,660 revenue → 4.83x ROAS (gross), 1.8x net of supply cost.",
                font_size=11, bold=True, color=FOREST, align=PP_ALIGN.CENTER)
    add_footer(s, PAGE)
slide_funnel()

# ===== SECTION 04 ===========================================
new_section_divider("04", "The product", "Live landing page · responsive · destination detail"); PAGE += 1

# ---- LANDING PAGE ARCHITECTURE
def slide_lp_arch():
    global PAGE; PAGE += 1
    s = new_content_slide("Landing page architecture", "Single-purpose page",
                          "One job: get the user to download the BB50 coupon · 11 sections, 1 conversion goal")
    sections = [
        ("Sticky nav", "Wordmark · Destinations · How it works · Blog · Coupon CTA · EN/IT toggle"),
        ("Hero", "Eyebrow · headline · sub · dual CTA · countdown timer · destination mosaic bg"),
        ("Value props", "4 pillars: curated · transparent · authentic · fast"),
        ("Destinations", "5 cards · price · nights · hover preview"),
        ("Comparison table", "BB&B vs OTAs (Booking, Expedia)"),
        ("How it works", "3 steps: pick · book · go"),
        ("Testimonials", "Quote carousel"),
        ("Social proof", "4-stat strip: travellers · destinations · save"),
        ("FAQ", "6 questions · accordion · tracked"),
        ("Lead form", "Name · email · departure · interest · consent · BB50 magnet"),
        ("Newsletter footer", "Tuesday weekly deal · separate from coupon flow"),
    ]
    y0 = 2.5
    for i, (name, desc) in enumerate(sections):
        col = i % 2
        row = i // 2
        x = 0.6 + col * 6.18
        yy = y0 + row * 0.7
        add_textbox(s, Inches(x), Inches(yy), Inches(0.3), Inches(0.3),
                    f"{i+1:02d}", font_size=10, bold=True, color=BURNT, font_name="Times-Bold")
        add_textbox(s, Inches(x + 0.5), Inches(yy), Inches(2), Inches(0.3),
                    name, font_size=11, bold=True, color=FOREST)
        add_textbox(s, Inches(x + 2.5), Inches(yy + 0.02), Inches(3.5), Inches(0.3),
                    desc, font_size=9, color=MUTED)
    add_textbox(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
                "Built on Next.js 14 (App Router) + TypeScript + Tailwind. Deployed on Vercel · sub-1s LCP · responsive.",
                font_size=10, italic=True, color=MUTED)
    add_footer(s, PAGE)
slide_lp_arch()

# ---- LIVE HOMEPAGE HERO
def slide_live_hero():
    global PAGE; PAGE += 1
    s = new_content_slide("Live · Homepage hero", "Above the fold",
                          "Dual CTA · countdown · destination mosaic · social proof eyebrow")
    add_image_below_title(s, os.path.join(SHOTS, "home-hero.png"),
                          top_in=2.4, max_w_in=11.5, max_h_in=4.4)
    add_footer(s, PAGE)
slide_live_hero()

# ---- LIVE FULL HOMEPAGE
def slide_live_full():
    global PAGE; PAGE += 1
    s = new_content_slide("Live · Full homepage", "All 11 sections of the landing flow",
                          "backpack-bottle.vercel.app · scroll-aware sticky CTA · Italian primary")
    img = os.path.join(SHOTS, "home.png")
    if os.path.exists(img):
        pi = PILImage.open(img)
        ar = pi.size[1] / pi.size[0]
        final_h = 4.5
        final_w = final_h / ar
        if final_w > 11:
            final_w = 11; final_h = final_w * ar
        left = (13.333 - final_w) / 2
        s.shapes.add_picture(img, Inches(left), Inches(2.4),
                             width=Inches(final_w), height=Inches(final_h))
    add_footer(s, PAGE)
slide_live_full()

# ---- LIVE MOBILE
def slide_live_mobile():
    global PAGE; PAGE += 1
    s = new_content_slide("Live · Mobile responsive", "375px viewport · iPhone-class",
                          "68% of Italian millennial sessions are mobile · same dataLayer events fire on every device")
    img = os.path.join(SHOTS, "home-mobile.png")
    if os.path.exists(img):
        pi = PILImage.open(img)
        ar = pi.size[1] / pi.size[0]
        final_h = 4.5
        final_w = final_h / ar
        left = (13.333 - final_w) / 2
        s.shapes.add_picture(img, Inches(left), Inches(2.4),
                             width=Inches(final_w), height=Inches(final_h))
    add_footer(s, PAGE)
slide_live_mobile()

# ---- LIVE DESTINATION DETAIL (Lisbon)
def slide_live_dest():
    global PAGE; PAGE += 1
    s = new_content_slide("Live · Destination detail", "/destinazioni/lisbon",
                          "Each city has its own page: 3-day itinerary, neighbourhoods, food, FAQ, weather, budget tip · IT + EN")
    add_image_below_title(s, os.path.join(SHOTS, "dest-lisbon.png"),
                          top_in=2.4, max_w_in=11.5, max_h_in=4.4)
    add_footer(s, PAGE)
slide_live_dest()

# ===== SECTION 05 ===========================================
new_section_divider("05", "Tracking architecture",
                    "GTM · GA4 · dataLayer · Marica's section"); PAGE += 1

# ---- THE FULL STACK VISUAL (the centerpiece)
def slide_stack():
    global PAGE; PAGE += 1
    s = new_content_slide("The full stack", "Six layers · 27 components",
                          "User journey: ad → site → form → coupon → email → booking")
    layers = [
        ("01 · ACQUISITION", "How users arrive", FOREST,
         ["Meta Ads (60%)", "Google Ads (40%)", "Organic SEO", "Newsletter"]),
        ("02 · ENGAGEMENT", "What they see", FOREST_DARK,
         ["Next.js 14 / Vercel", "IT/EN toggle", "MS Clarity", "Exit-intent", "WhatsApp"]),
        ("03 · CONVERSION", "How they act", BURNT,
         ["Lead form", "BB50 €50", "Countdown", "Sticky CTA", "/grazie"]),
        ("04 · RETENTION", "How they come back", BURNT_DARK,
         ["Brevo ESP", "Automation #23", "UTM persistence", "Meta CAPI"]),
        ("05 · MEASUREMENT", "How we know it worked", FOREST,
         ["GTM-KSN24TZ2", "GA4 G-FHBQE8QZM9", "Meta Pixel", "Looker Studio", "Consent v2"]),
        ("06 · INFRA", "Where it lives", MUTED,
         ["Vercel", "GitHub (public)", "Next.js metadata", "GDPR + rate-limit"]),
    ]
    y0 = 2.4
    band_h = 0.7
    total_w = 12.13
    label_w = 2.0
    for i, (label, sub, accent, tools) in enumerate(layers):
        y = y0 + i * (band_h + 0.05)
        add_rect(s, Inches(0.6), Inches(y), Inches(total_w), Inches(band_h),
                 CREAM_DARK, line_color=accent)
        add_rect(s, Inches(0.6), Inches(y), Inches(label_w), Inches(band_h), accent)
        add_textbox(s, Inches(0.7), Inches(y + 0.08), Inches(label_w - 0.2), Inches(0.3),
                    label, font_size=10, bold=True, color=CREAM)
        add_textbox(s, Inches(0.7), Inches(y + 0.32), Inches(label_w - 0.2), Inches(0.3),
                    sub, font_size=8, color=CREAM, italic=True)
        x = 0.6 + label_w + 0.15
        pill_h = 0.32
        pill_y = y + (band_h - pill_h) / 2
        for tool in tools:
            tw = max(1.4, len(tool) * 0.085)
            if x + tw > 0.6 + total_w - 0.1: break
            add_rect(s, Inches(x), Inches(pill_y), Inches(tw), Inches(pill_h),
                     WHITE, line_color=accent)
            add_textbox(s, Inches(x + 0.08), Inches(pill_y + 0.04), Inches(tw - 0.15), Inches(0.25),
                        tool, font_size=8.5, color=FOREST, align=PP_ALIGN.CENTER)
            x += tw + 0.1
    add_textbox(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
                "All 27 components shipped. Single source of truth: window.dataLayer · all consent-gated.",
                font_size=9.5, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, PAGE)
slide_stack()

# ---- TRACKING DATAFLOW
def slide_dataflow():
    global PAGE; PAGE += 1
    s = new_content_slide("Tracking architecture", "One container. Every tag. Consent-aware.",
                          "How a single user click becomes 4 events in 4 different systems")
    # Diagram (text-based)
    add_rect(s, Inches(2), Inches(2.5), Inches(9.5), Inches(4.0), FOREST_DARK)
    code = [
        ("USER ACTION", "Click 'Get your €50' / Submit form / Download PDF", BURNT),
        ("│", "", CREAM),
        ("▼", "", BURNT),
        ("dataLayer.push({event: 'coupon_download', value: 50, currency: 'EUR'})",
         "single source of truth · always fires · consent-independent", CREAM),
        ("│", "", CREAM),
        ("▼", "GTM container GTM-KSN24TZ2 listens", BURNT),
        ("┌────────────┬─────────────┬───────────────┬──────────────┐", "", CREAM),
        ("│  GA4       │  Meta Pixel │  Meta CAPI    │  Google Ads  │", "", CREAM),
        ("│  events    │  Lead+CR    │  server-side  │  conversion  │", "", CREAM),
        ("└────────────┴─────────────┴───────────────┴──────────────┘", "", CREAM),
        ("ALL GATED BY CONSENT MODE V2 ·  500ms wait_for_update", "", BURNT),
    ]
    y = 2.7
    for i, (line, sub, c) in enumerate(code):
        is_special = line in ["│", "▼"]
        font_size = 14 if is_special else 9.5
        font = "Helvetica" if is_special else "Courier"
        bold = i in [0, 5, 10]
        add_textbox(s, Inches(2), Inches(y + i*0.32), Inches(9.5), Inches(0.32),
                    line, font_size=font_size, bold=bold, color=c, font_name=font,
                    align=PP_ALIGN.CENTER)
    add_footer(s, PAGE)
slide_dataflow()

# ---- DATALAYER SCHEMA TABLE
def slide_datalayer():
    global PAGE; PAGE += 1
    s = new_content_slide("dataLayer event schema", "31 events · 11 critical · 3 Key Events",
                          "Every interaction on the site becomes a structured dataLayer push")
    events = [
        ("page_view", "Every route", "Acquisition baseline", False),
        ("cta_click", "Hero / Nav / Sticky CTA", "Which CTA placement converts", False),
        ("form_view", "Form 40% in viewport", "Funnel impression", False),
        ("form_start", "First field focus", "Intent signal", False),
        ("form_submit", "Server confirms ok=true", "Lead conversion", True),
        ("coupon_download", "/grazie + PDF link click", "€50 value goal", True),
        ("thank_you_view", "/grazie page mount", "Server-confirm", True),
        ("destination_card_click", "City card click", "destination_interest custom dim", False),
        ("scroll_75", "75% scroll depth", "Engagement quality", False),
        ("rage_click", "≥3 same coords <1s", "UX defect detector", False),
        ("consent_choice", "Accept/Reject click", "GDPR audit trail (NEW)", False),
        ("consent_update", "Cookie banner choice", "Consent Mode v2 state", False),
    ]
    add_rect(s, Inches(0.6), Inches(2.45), Inches(12.13), Inches(0.4), FOREST)
    cols = [Inches(2.7), Inches(3.5), Inches(4.5), Inches(1.5)]
    xs = [Inches(0.6)]
    for w in cols[:-1]: xs.append(xs[-1] + w)
    headers = ["EVENT", "TRIGGER", "PURPOSE", "KEY EVENT?"]
    for i, h in enumerate(headers):
        add_textbox(s, xs[i] + Inches(0.1), Inches(2.55), cols[i], Inches(0.3),
                    h, font_size=9, bold=True, color=CREAM)
    y = 2.85
    for i, (name, trigger, purpose, key) in enumerate(events):
        bg = WHITE if i % 2 == 0 else CREAM_DARK
        add_rect(s, Inches(0.6), Inches(y + i*0.34), Inches(12.13), Inches(0.34), bg)
        add_textbox(s, xs[0] + Inches(0.1), Inches(y + i*0.34 + 0.05), cols[0], Inches(0.3),
                    name, font_size=9, bold=True, color=BURNT_DARK, font_name="Courier")
        add_textbox(s, xs[1] + Inches(0.1), Inches(y + i*0.34 + 0.05), cols[1], Inches(0.3),
                    trigger, font_size=9, color=INK)
        add_textbox(s, xs[2] + Inches(0.1), Inches(y + i*0.34 + 0.05), cols[2], Inches(0.3),
                    purpose, font_size=9, color=INK)
        add_textbox(s, xs[3] + Inches(0.1), Inches(y + i*0.34 + 0.05), cols[3], Inches(0.3),
                    "★ YES" if key else "—",
                    font_size=9, bold=key, color=BURNT if key else MUTED)

    add_textbox(s, Inches(0.6), Inches(7.1), Inches(12), Inches(0.3),
                "+ 19 more events instrumented · full schema at backpack-bottle.vercel.app/measurement",
                font_size=9, italic=True, color=MUTED)
    add_footer(s, PAGE)
slide_datalayer()

# ---- LIVE GTM (mockup)
def slide_gtm_live():
    global PAGE; PAGE += 1
    s = new_content_slide("Live · GTM workspace", "GTM-KSN24TZ2 · published v5",
                          "17 tags · 22 triggers · 17 variables · imported via JSON, not click-by-click")
    add_image_below_title(s, os.path.join(SHOTS, "gtm.png"),
                          top_in=2.4, max_w_in=11.5, max_h_in=4.4)
    add_footer(s, PAGE)
slide_gtm_live()

# ---- LIVE GA4 (mockup)
def slide_ga4_live():
    global PAGE; PAGE += 1
    s = new_content_slide("Live · GA4 Realtime", "G-FHBQE8QZM9 · 5 custom dims · 3 Key Events",
                          "Events flow within 30 seconds of any interaction · debug-mode confirmed")
    add_image_below_title(s, os.path.join(SHOTS, "ga4.png"),
                          top_in=2.4, max_w_in=11.5, max_h_in=4.4)
    add_footer(s, PAGE)
slide_ga4_live()

# ---- LIVE /MEASUREMENT
def slide_measurement_live():
    global PAGE; PAGE += 1
    s = new_content_slide("Live · /measurement", "Every event documented on the site",
                          "KPI tree · dataLayer schema · custom dims · Consent Mode v2 audit · attribution model")
    add_image_below_title(s, os.path.join(SHOTS, "measurement.png"),
                          top_in=2.4, max_w_in=11.5, max_h_in=4.4)
    add_footer(s, PAGE)
slide_measurement_live()

# ===== SECTION 06 ===========================================
new_section_divider("06", "Compliance",
                    "Consent Mode v2 · cookie tracking · GDPR · Alberto's section"); PAGE += 1

# ---- CONSENT MODE V2
def slide_consent():
    global PAGE; PAGE += 1
    s = new_content_slide("Consent Mode v2", "Default-denied. Privacy-first. Audit-ready.",
                          "Required for any site serving EU traffic · compliant with GDPR + ePrivacy")
    # Left side — explanation
    add_textbox(s, Inches(0.6), Inches(2.45), Inches(2.5), Inches(0.3),
                "HOW IT WORKS", font_size=10, bold=True, color=BURNT)
    bullets = [
        "Before GTM loads, gtag('consent','default',...) sets all four signals to denied.",
        "Cookie banner asks the user. On Accept → consent_update with all granted.",
        "On Reject → stays denied; GA4 receives consent-modeled pings (no cookies).",
        "500ms wait_for_update prevents tags from firing before consent is set.",
        "Audit trail: every consent_update + consent_choice event stored in GA4.",
    ]
    for i, b in enumerate(bullets):
        add_textbox(s, Inches(0.85), Inches(2.85 + i*0.55), Inches(0.3), Inches(0.3),
                    "·", font_size=20, bold=True, color=BURNT)
        add_textbox(s, Inches(1.1), Inches(2.92 + i*0.55), Inches(5.7), Inches(0.7),
                    b, font_size=10, color=INK, line_spacing=1.3)

    # Right side — code
    add_rect(s, Inches(7.2), Inches(2.4), Inches(5.7), Inches(4.4), FOREST_DARK)
    add_textbox(s, Inches(7.4), Inches(2.5), Inches(5.3), Inches(0.3),
                "components/GTM.tsx", font_size=8.5, color=BURNT, font_name="Courier")
    code_lines = [
        "gtag('consent', 'default', {",
        "  ad_storage:        'denied',",
        "  analytics_storage: 'denied',",
        "  ad_user_data:      'denied',",
        "  ad_personalization:'denied',",
        "  wait_for_update: 500,",
        "});",
        "",
        "// on user click in CookieBanner:",
        "gtag('consent', 'update', {",
        "  ad_storage:        'granted',",
        "  analytics_storage: 'granted',",
        "  ad_user_data:      'granted',",
        "  ad_personalization:'granted',",
        "});",
    ]
    for i, line in enumerate(code_lines):
        add_textbox(s, Inches(7.4), Inches(2.85 + i*0.25), Inches(5.3), Inches(0.25),
                    line, font_size=9, color=CREAM, font_name="Courier")
    add_footer(s, PAGE)
slide_consent()

# ---- COOKIE CONSENT TRACKING (NEW)
def slide_cookie_tracking():
    global PAGE; PAGE += 1
    s = new_content_slide("Cookie consent — measured, not just enforced",
                          "Every Accept and Reject click is a tracked event",
                          "consent_choice fires alongside consent_update · dashboard shows Accept-rate vs Reject-rate")
    # Left — what fires
    add_textbox(s, Inches(0.6), Inches(2.45), Inches(2.5), Inches(0.3),
                "WHAT WE TRACK", font_size=10, bold=True, color=BURNT)
    add_multiline(s, Inches(0.6), Inches(2.85), Inches(6.0), Inches(2.5), [
        "1. consent_update — sets the four Consent Mode v2 signals",
        "    (Google's spec · required for compliant tag firing)",
        "",
        "2. consent_choice — discrete event with two params:",
        "    consent_choice = 'accept' | 'reject'",
        "    consent_locale = 'it' | 'en'",
        "",
        "GTM forwards consent_choice to GA4 as a key event.",
        "Looker Studio shows real-time Accept-rate vs Reject-rate.",
    ], font_size=10, color=INK, line_spacing=1.3)

    # Numbers expected
    add_textbox(s, Inches(0.6), Inches(5.7), Inches(2.5), Inches(0.3),
                "INDUSTRY BENCHMARKS", font_size=10, bold=True, color=BURNT)
    nums = [
        ("Accept rate (EU)", "~62%"),
        ("Reject rate (EU)", "~28%"),
        ("Banner-ignored", "~10%"),
    ]
    for i, (k, v) in enumerate(nums):
        add_textbox(s, Inches(0.6), Inches(6.05 + i*0.3), Inches(3), Inches(0.25),
                    k, font_size=10, color=MUTED)
        add_textbox(s, Inches(3.7), Inches(6.05 + i*0.3), Inches(2), Inches(0.25),
                    v, font_size=10, bold=True, color=FOREST, font_name="Courier")

    # Right — code
    add_rect(s, Inches(7.2), Inches(2.4), Inches(5.7), Inches(4.4), FOREST_DARK)
    add_textbox(s, Inches(7.4), Inches(2.5), Inches(5.3), Inches(0.3),
                "components/CookieBanner.tsx", font_size=8.5, color=BURNT, font_name="Courier")
    code = [
        "function decide(granted: boolean) {",
        "  localStorage.setItem('bb_consent',",
        "    granted ? 'granted' : 'denied');",
        "  pushConsent(granted);",
        "",
        "  // Discrete tracked event:",
        "  dataLayer.push({",
        "    event: 'consent_choice',",
        "    consent_choice: granted",
        "      ? 'accept' : 'reject',",
        "    consent_locale: locale,",
        "  });",
        "",
        "  setShow(false);",
        "}",
    ]
    for i, line in enumerate(code):
        add_textbox(s, Inches(7.4), Inches(2.85 + i*0.25), Inches(5.3), Inches(0.25),
                    line, font_size=9, color=CREAM, font_name="Courier")
    add_footer(s, PAGE)
slide_cookie_tracking()

# ---- LIVE PRIVACY
def slide_privacy_live():
    global PAGE; PAGE += 1
    s = new_content_slide("Live · /privacy", "GDPR + cookie + coupon T&Cs",
                          "Seven sections · Italian + English · audit-ready")
    add_image_below_title(s, os.path.join(SHOTS, "privacy.png"),
                          top_in=2.4, max_w_in=11.5, max_h_in=4.4)
    add_footer(s, PAGE)
slide_privacy_live()

# ===== SECTION 07 ===========================================
new_section_divider("07", "Email & automation",
                    "Brevo · welcome series · Ceci + Marilu's section"); PAGE += 1

# ---- FORM → API → BREVO PIPELINE
def slide_pipeline():
    global PAGE; PAGE += 1
    s = new_content_slide("Lead pipeline", "Form → API → Brevo (parallel fan-out)",
                          "Sub-200ms server response · 3 systems updated in parallel · then automation fires")
    boxes = [
        (1.0, "Form submit", "User clicks 'Get €50'\non backpack-bottle.vercel.app", FOREST),
        (4.0, "POST /api/lead", "Server validates email\nrate-limits 5/min/IP", BURNT),
        (7.0, "Promise.all", "3 fan-outs in parallel:", FOREST_DARK),
        (10.0, "Returns ok:true", "Client redirects to /grazie\nfires Key Events", BURNT_DARK),
    ]
    y = 2.7
    for i, (x, lbl, sub, col) in enumerate(boxes):
        add_rect(s, Inches(x), Inches(y), Inches(2.6), Inches(1.5), col)
        add_textbox(s, Inches(x), Inches(y + 0.2), Inches(2.6), Inches(0.4),
                    lbl, font_size=12, bold=True, color=CREAM, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(x + 0.15), Inches(y + 0.7), Inches(2.3), Inches(0.7),
                    sub, font_size=9, color=CREAM, align=PP_ALIGN.CENTER, italic=True)
        if i < 3:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     Inches(x + 2.6), Inches(y + 0.6),
                                     Inches(0.4), Inches(0.3))
            ar.fill.solid(); ar.fill.fore_color.rgb = BURNT
            ar.line.fill.background()

    # Fan-out detail
    add_rect(s, Inches(0.6), Inches(4.7), Inches(12.13), Inches(2.0), CREAM_DARK)
    add_textbox(s, Inches(0.85), Inches(4.85), Inches(2.5), Inches(0.3),
                "PARALLEL FAN-OUT (Promise.all)", font_size=10, bold=True, color=BURNT)
    fanout = [
        ("Brevo API", "POST /v3/contacts → list 13 with 7 attributes\nupdateEnabled:true (idempotent)"),
        ("Meta CAPI", "POST graph.facebook.com/v18.0/{pixel}/events\nLead event · SHA-256 email · fbp/fbc cookies"),
        ("Webhook (optional)", "POST $LEAD_WEBHOOK_URL\nfor Zapier / Sheets / Slack alerts"),
    ]
    for i, (name, desc) in enumerate(fanout):
        x = 0.85 + i * 4.0
        add_textbox(s, Inches(x), Inches(5.25), Inches(3.8), Inches(0.3),
                    name, font_size=11, bold=True, color=FOREST)
        add_textbox(s, Inches(x), Inches(5.55), Inches(3.8), Inches(1.0),
                    desc, font_size=9, color=INK, font_name="Courier", line_spacing=1.3)
    add_footer(s, PAGE)
slide_pipeline()

# ---- LIVE BREVO (mockup)
def slide_brevo_live():
    global PAGE; PAGE += 1
    s = new_content_slide("Live · Brevo dashboard", "List 'backpackandbottle' · ID 13",
                          "Master list · all form submissions · trigger for Automation #23 · 40 templates ready")
    add_image_below_title(s, os.path.join(SHOTS, "brevo.png"),
                          top_in=2.4, max_w_in=11.5, max_h_in=4.4)
    add_footer(s, PAGE)
slide_brevo_live()

# ---- 8-EMAIL WELCOME SERIES
def slide_welcome_series():
    global PAGE; PAGE += 1
    s = new_content_slide("8-email welcome series", "Automation #23 · BB50 Welcome Journey",
                          "Triggers when contact lands in list 13 · personalised by destination · IT/EN by locale attribute")
    emails = [
        ("01", "T+0", "Coupon delivery", "Il tuo coupon BB50 è arrivato 🎒🍷"),
        ("02", "T+1d", "Brand intro", "Perché Backpack & Bottle?"),
        ("03", "T+3d", "Destination deep-dive", "[Personalised: Lisbona / Roma / ...]"),
        ("04", "T+7d", "Booking psychology", "Quando prenotare il volo (dati 2026)"),
        ("05", "T+14d", "Social proof", "Le 3 città più richieste questa settimana"),
        ("06", "T+30d", "Soft reminder", "Hai ancora 60 giorni per usare BB50"),
        ("07", "T+60d", "Expiry warning", "Il tuo coupon BB50 scade tra 30 giorni"),
        ("08", "T+83d", "Final urgency", "Ultimi 7 giorni per il tuo €50"),
    ]
    y = 2.45
    for i, (num, when, kind, subject) in enumerate(emails):
        col = i % 2
        row = i // 2
        x = 0.6 + col * 6.18
        yy = y + row * 1.1
        add_rect(s, Inches(x), Inches(yy), Inches(5.95), Inches(0.95), WHITE,
                 line_color=CREAM_DARK)
        add_rect(s, Inches(x), Inches(yy), Inches(0.85), Inches(0.95), BURNT)
        add_textbox(s, Inches(x), Inches(yy + 0.1), Inches(0.85), Inches(0.3),
                    num, font_size=18, bold=True, color=CREAM, align=PP_ALIGN.CENTER,
                    font_name="Times-Bold")
        add_textbox(s, Inches(x), Inches(yy + 0.5), Inches(0.85), Inches(0.3),
                    when, font_size=9, color=CREAM, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(x + 1.05), Inches(yy + 0.1), Inches(4.7), Inches(0.3),
                    kind, font_size=10, bold=True, color=BURNT_DARK)
        add_textbox(s, Inches(x + 1.05), Inches(yy + 0.4), Inches(4.7), Inches(0.5),
                    subject, font_size=10, color=INK, italic=True)
    add_footer(s, PAGE)
slide_welcome_series()

# ===== SECTION 08 ===========================================
new_section_divider("08", "Media plan",
                    "€2,000 · 60% Meta · 40% Google · 6-week ramp"); PAGE += 1

# ---- MEDIA PLAN OVERVIEW
def slide_media():
    global PAGE; PAGE += 1
    s = new_content_slide("Media plan overview", "Where the €2,000 goes",
                          "Two channels · one job each · Meta drives discovery, Google captures intent")
    # Big split
    add_rect(s, Inches(0.6), Inches(2.5), Inches(7.4), Inches(4.4), FOREST)
    add_textbox(s, Inches(0.6), Inches(2.75), Inches(7.4), Inches(0.4),
                "META ADS · 60% · €1,200", font_size=11, bold=True, color=BURNT,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(3.2), Inches(7.4), Inches(0.7),
                "Discovery + retargeting", font_size=24, color=CREAM,
                font_name="Times-Bold", align=PP_ALIGN.CENTER)
    rows_m = [
        ("Awareness", "€300", "Carousel ads · broad reach · 25–39 IT"),
        ("Lead Gen", "€600", "Conversion-optimised · interest targeting"),
        ("Retargeting", "€300", "DPA · 30-day window · /destinazioni viewers"),
    ]
    for i, (lbl, amt, desc) in enumerate(rows_m):
        y = 4.1 + i*0.85
        add_textbox(s, Inches(0.9), Inches(y), Inches(2), Inches(0.3),
                    lbl, font_size=11, bold=True, color=BURNT)
        add_textbox(s, Inches(2.9), Inches(y), Inches(1.2), Inches(0.3),
                    amt, font_size=11, color=CREAM, font_name="Courier", bold=True)
        add_textbox(s, Inches(0.9), Inches(y + 0.32), Inches(7), Inches(0.4),
                    desc, font_size=10, color=CREAM, italic=True)

    add_rect(s, Inches(8.2), Inches(2.5), Inches(4.5), Inches(4.4), CREAM_DARK)
    add_textbox(s, Inches(8.2), Inches(2.75), Inches(4.5), Inches(0.4),
                "GOOGLE ADS · 40% · €800", font_size=11, bold=True, color=BURNT,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(8.2), Inches(3.2), Inches(4.5), Inches(0.7),
                "Active intent capture", font_size=20, color=FOREST,
                font_name="Times-Bold", align=PP_ALIGN.CENTER)
    rows_g = [
        ("Branded", "€200", "Defend 'backpack & bottle' searches"),
        ("Generic", "€480", "'weekend a Lisbona', 'voli + hotel'"),
        ("Competitor", "€120", "'volagratis', 'edreams' bidding"),
    ]
    for i, (lbl, amt, desc) in enumerate(rows_g):
        y = 4.1 + i*0.85
        add_textbox(s, Inches(8.4), Inches(y), Inches(1.8), Inches(0.3),
                    lbl, font_size=11, bold=True, color=BURNT_DARK)
        add_textbox(s, Inches(10.2), Inches(y), Inches(1.0), Inches(0.3),
                    amt, font_size=11, color=FOREST, font_name="Courier", bold=True)
        add_textbox(s, Inches(8.4), Inches(y + 0.32), Inches(4.0), Inches(0.4),
                    desc, font_size=9, color=INK, italic=True)
    add_footer(s, PAGE)
slide_media()

# ---- BUDGET PACING
def slide_pacing():
    global PAGE; PAGE += 1
    s = new_content_slide("Budget pacing", "€2,000 across 6 weeks",
                          "Front-loaded for discovery · scaled into peak booking-decision window")
    weeks = [
        ("W1", "May 5–11", 200, 100, 30),  # meta, google, expected leads
        ("W2", "May 12–18", 240, 140, 51),
        ("W3", "May 19–25", 240, 160, 68),
        ("W4", "May 26–Jun 1", 220, 160, 64),
        ("W5", "Jun 2–8", 180, 140, 51),
        ("W6", "Jun 9–15", 120, 100, 36),
    ]
    # Header row
    add_rect(s, Inches(0.6), Inches(2.5), Inches(12.13), Inches(0.4), FOREST)
    headers = ["WEEK", "DATES", "META", "GOOGLE", "TOTAL", "PROJECTED LEADS"]
    cw = [Inches(1.0), Inches(2.5), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.63)]
    xs = [Inches(0.6)]
    for w in cw[:-1]: xs.append(xs[-1] + w)
    for i, h in enumerate(headers):
        add_textbox(s, xs[i] + Inches(0.1), Inches(2.6), cw[i], Inches(0.3),
                    h, font_size=9, bold=True, color=CREAM)
    y = 2.9
    total_meta, total_google, total_leads = 0, 0, 0
    for i, (wk, dates, meta, google, leads) in enumerate(weeks):
        bg = WHITE if i % 2 == 0 else CREAM_DARK
        add_rect(s, Inches(0.6), Inches(y + i*0.5), Inches(12.13), Inches(0.5), bg)
        add_textbox(s, xs[0] + Inches(0.1), Inches(y + i*0.5 + 0.1), cw[0], Inches(0.3),
                    wk, font_size=11, bold=True, color=BURNT)
        add_textbox(s, xs[1] + Inches(0.1), Inches(y + i*0.5 + 0.13), cw[1], Inches(0.3),
                    dates, font_size=10, color=INK)
        add_textbox(s, xs[2] + Inches(0.1), Inches(y + i*0.5 + 0.13), cw[2], Inches(0.3),
                    f"€{meta}", font_size=10, color=FOREST, font_name="Courier")
        add_textbox(s, xs[3] + Inches(0.1), Inches(y + i*0.5 + 0.13), cw[3], Inches(0.3),
                    f"€{google}", font_size=10, color=FOREST, font_name="Courier")
        add_textbox(s, xs[4] + Inches(0.1), Inches(y + i*0.5 + 0.13), cw[4], Inches(0.3),
                    f"€{meta+google}", font_size=10, bold=True, color=FOREST, font_name="Courier")
        # Visual bar for leads
        bar_w = leads / 70 * 2.4
        add_rect(s, xs[5] + Inches(0.1), Inches(y + i*0.5 + 0.18), Inches(bar_w), Inches(0.18), BURNT)
        add_textbox(s, xs[5] + Inches(bar_w + 0.15), Inches(y + i*0.5 + 0.13), Inches(1.0), Inches(0.3),
                    f"{leads}", font_size=10, bold=True, color=BURNT_DARK)
        total_meta += meta; total_google += google; total_leads += leads

    # Total row
    yt = y + 6*0.5
    add_rect(s, Inches(0.6), Inches(yt), Inches(12.13), Inches(0.5), FOREST_DARK)
    add_textbox(s, xs[0] + Inches(0.1), Inches(yt + 0.1), cw[0], Inches(0.3),
                "TOTAL", font_size=11, bold=True, color=BURNT)
    add_textbox(s, xs[2] + Inches(0.1), Inches(yt + 0.13), cw[2], Inches(0.3),
                f"€{total_meta}", font_size=10, bold=True, color=CREAM, font_name="Courier")
    add_textbox(s, xs[3] + Inches(0.1), Inches(yt + 0.13), cw[3], Inches(0.3),
                f"€{total_google}", font_size=10, bold=True, color=CREAM, font_name="Courier")
    add_textbox(s, xs[4] + Inches(0.1), Inches(yt + 0.13), cw[4], Inches(0.3),
                f"€{total_meta + total_google}", font_size=10, bold=True, color=CREAM, font_name="Courier")
    add_textbox(s, xs[5] + Inches(0.1), Inches(yt + 0.13), cw[5], Inches(0.3),
                f"{total_leads} leads",
                font_size=10, bold=True, color=BURNT)
    add_footer(s, PAGE)
slide_pacing()

# ===== SECTION 09 ===========================================
new_section_divider("09", "KPIs & measurement",
                    "Funnel targets · Looker Studio · Amelia's section"); PAGE += 1

# ---- KPI TREE
def slide_kpis():
    global PAGE; PAGE += 1
    s = new_content_slide("KPI tree & targets", "8 metrics · ladder from sessions to bookings",
                          "Targets derived from €2,000 spend × 60/40 channel split × Italian travel benchmarks")
    rows = [
        ("Top of funnel", "Sessions", "5,000", "GA4 Acquisition"),
        ("Mid funnel", "Form views", "1,500 (30%)", "form_view event"),
        ("Mid funnel", "Form starts", "750 (50%)", "form_start event"),
        ("Bottom funnel", "Leads (form_submit)", "300 (40%)", "★ Key Event"),
        ("Bottom funnel", "Coupon downloads", "400", "★ Key Event"),
        ("Post-campaign", "Bookings (offline)", "30", "Brevo → GA4 import"),
        ("Efficiency", "CPL", "€6.67", "Spend ÷ leads"),
        ("Efficiency", "ROAS (proj)", "1.8x", "AOV €322 × bookings ÷ spend"),
    ]
    add_rect(s, Inches(0.6), Inches(2.5), Inches(12.13), Inches(0.4), FOREST)
    headers = ["STAGE", "METRIC", "TARGET (6w)", "SOURCE"]
    cw = [Inches(2.6), Inches(3.5), Inches(2.5), Inches(3.53)]
    xs = [Inches(0.6)]
    for w in cw[:-1]: xs.append(xs[-1] + w)
    for i, h in enumerate(headers):
        add_textbox(s, xs[i] + Inches(0.1), Inches(2.6), cw[i], Inches(0.3),
                    h, font_size=9, bold=True, color=CREAM)
    y = 2.9
    for i, (stage, metric, target, source) in enumerate(rows):
        bg = WHITE if i % 2 == 0 else CREAM_DARK
        add_rect(s, Inches(0.6), Inches(y + i*0.45), Inches(12.13), Inches(0.45), bg)
        add_textbox(s, xs[0] + Inches(0.1), Inches(y + i*0.45 + 0.1), cw[0], Inches(0.3),
                    stage, font_size=10, color=MUTED)
        add_textbox(s, xs[1] + Inches(0.1), Inches(y + i*0.45 + 0.1), cw[1], Inches(0.3),
                    metric, font_size=11, bold=True, color=FOREST)
        add_textbox(s, xs[2] + Inches(0.1), Inches(y + i*0.45 + 0.1), cw[2], Inches(0.3),
                    target, font_size=12, bold=True, color=BURNT_DARK, font_name="Times-Bold")
        add_textbox(s, xs[3] + Inches(0.1), Inches(y + i*0.45 + 0.1), cw[3], Inches(0.3),
                    source, font_size=9, color=MUTED, italic=True)
    add_footer(s, PAGE)
slide_kpis()

# ---- LOOKER STUDIO PLACEHOLDER (replaces whacky mock)
def slide_looker_placeholder():
    global PAGE; PAGE += 1
    s = new_content_slide("Looker Studio dashboard", "Single-page exec view · refreshes every 12 hours",
                          "What we look at every morning during the campaign")
    add_rect(s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(3.2), CREAM_DARK,
             line_color=BURNT, line_width=Pt(2))
    add_textbox(s, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.5),
                "[ Live dashboard screenshot to be inserted on launch day ]",
                font_size=14, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.5),
                "Built in Looker Studio · 6 panels · refreshes every 12 hours",
                font_size=11, color=FOREST, align=PP_ALIGN.CENTER)
    panels = [
        ("ACQUISITION", "Sessions, users by source/medium · weekly trend"),
        ("FUNNEL", "Sessions → form_submit → coupon · drop-off %"),
        ("COST", "Spend, CPM, CPC, CTR, CPL by channel"),
        ("AUDIENCE", "Geo (IT regions), device, locale (it/en)"),
        ("ENGAGEMENT", "Avg session, scroll_75%, engaged %"),
        ("COUPON", "Downloads by destination_interest custom dim"),
    ]
    y = 5.95
    for i, (lbl, desc) in enumerate(panels):
        col = i % 3
        row = i // 3
        x = 0.6 + col * 4.18
        yy = y + row * 0.55
        add_textbox(s, Inches(x), Inches(yy), Inches(1.5), Inches(0.3),
                    lbl, font_size=9, bold=True, color=BURNT)
        add_textbox(s, Inches(x + 1.55), Inches(yy + 0.02), Inches(2.5), Inches(0.3),
                    desc, font_size=8, color=INK)
    add_footer(s, PAGE)
slide_looker_placeholder()

# ===== SECTION 10 ===========================================
new_section_divider("10", "Optimization",
                    "A/B test · competitors · accessibility"); PAGE += 1

# ---- A/B TEST
def slide_ab():
    global PAGE; PAGE += 1
    s = new_content_slide("A/B test plan", "Day-1 hypothesis · hero CTA copy",
                          "Test what matters most: the headline that gets the click")
    add_textbox(s, Inches(0.6), Inches(2.4), Inches(2.2), Inches(0.3),
                "HYPOTHESIS", font_size=10, bold=True, color=BURNT)
    add_textbox(s, Inches(0.6), Inches(2.7), Inches(12), Inches(1.0),
                "An emotional hero CTA ('Discover your weekend') will out-convert a transactional one "
                "('Get your €50 coupon') for cold awareness audiences. Italian millennials respond to "
                "aspiration over discount when not yet brand-aware.",
                font_size=11, color=INK, line_spacing=1.4)

    add_rect(s, Inches(0.6), Inches(4.0), Inches(6.0), Inches(1.4), CREAM_DARK)
    add_textbox(s, Inches(0.85), Inches(4.1), Inches(5.5), Inches(0.3),
                "VARIANT A · CONTROL", font_size=9, bold=True, color=BURNT)
    add_textbox(s, Inches(0.85), Inches(4.4), Inches(5.5), Inches(0.5),
                '"Get your €50 coupon →"', font_size=15, bold=True, color=FOREST,
                font_name="Times-Bold")
    add_textbox(s, Inches(0.85), Inches(4.95), Inches(5.5), Inches(0.4),
                "Transactional · current production copy", font_size=9, color=MUTED)

    add_rect(s, Inches(6.83), Inches(4.0), Inches(6.0), Inches(1.4), CREAM_DARK)
    add_textbox(s, Inches(7.08), Inches(4.1), Inches(5.5), Inches(0.3),
                "VARIANT B · CHALLENGER", font_size=9, bold=True, color=BURNT)
    add_textbox(s, Inches(7.08), Inches(4.4), Inches(5.5), Inches(0.5),
                '"Discover your next weekend →"', font_size=15, bold=True, color=FOREST,
                font_name="Times-Bold")
    add_textbox(s, Inches(7.08), Inches(4.95), Inches(5.5), Inches(0.4),
                "Emotional · for cold awareness audiences", font_size=9, color=MUTED)

    add_rect(s, Inches(0.6), Inches(5.6), Inches(12.13), Inches(1.4), FOREST)
    add_textbox(s, Inches(0.85), Inches(5.7), Inches(3), Inches(0.3),
                "DECISION CRITERIA", font_size=10, bold=True, color=BURNT)
    rows = [
        ("Primary metric", "form_submit conversion rate"),
        ("Sample size", "~750 sessions per variant (3% baseline · 20% MDE · 95% confidence)"),
        ("Test duration", "10–14 days (covers weekday + weekend behaviour)"),
        ("Decision rule", "Ship winner if uplift ≥15% with p<0.05; else keep control"),
    ]
    for i, (k, v) in enumerate(rows):
        add_textbox(s, Inches(0.85), Inches(6.0 + i*0.23), Inches(2.5), Inches(0.22),
                    k, font_size=9, bold=True, color=CREAM)
        add_textbox(s, Inches(3.45), Inches(6.0 + i*0.23), Inches(9.3), Inches(0.22),
                    v, font_size=9, color=CREAM)
    add_footer(s, PAGE)
slide_ab()

# ---- COMPETITOR
def slide_competitors():
    global PAGE; PAGE += 1
    s = new_content_slide("Competitor benchmark", "Italian travel sector · CPL, CTR, positioning",
                          "What we're up against · why our targets are realistic")
    headers = ["COMPETITOR", "POSITIONING", "CPL", "CTR", "LEAD MAGNET", "OUR EDGE"]
    col_widths = [Inches(2.0), Inches(2.5), Inches(1.0), Inches(1.0), Inches(2.5), Inches(2.5)]
    xs = [Inches(0.6)]
    for w in col_widths[:-1]: xs.append(xs[-1] + w)
    add_rect(s, Inches(0.6), Inches(2.5), Inches(11.5), Inches(0.4), FOREST)
    for i, h in enumerate(headers):
        add_textbox(s, xs[i] + Inches(0.1), Inches(2.6), col_widths[i], Inches(0.3),
                    h, font_size=9, bold=True, color=CREAM)
    rows = [
        ("Volagratis", "Cheapest flights, no curation", "€11.20", "1.8%",
         "Newsletter signup", "We curate; they don't"),
        ("eDreams Prime", "Membership flight discounts", "€9.50", "2.0%",
         "Free Prime trial", "No subscription friction"),
        ("Lastminute.com", "Last-minute travel deals", "€8.30", "2.4%",
         "App download", "We're early-booking, not last-minute"),
        ("BACKPACK & BOTTLE", "Curated · transparent · early", "€6.67",
         "2.7% (proj)", "BB50 €50 coupon", "First IT to combine all 3"),
    ]
    for r, row in enumerate(rows):
        yy = 2.9 + r*0.55
        bg = WHITE if r % 2 == 0 and r < 3 else (CREAM_DARK if r < 3 else BURNT)
        add_rect(s, Inches(0.6), Inches(yy), Inches(11.5), Inches(0.55), bg)
        for i, val in enumerate(row):
            text_color = CREAM if r == 3 else INK
            bold = (i == 0) or (r == 3)
            add_textbox(s, xs[i] + Inches(0.1), Inches(yy + 0.16),
                        col_widths[i], Inches(0.3), val,
                        font_size=10, bold=bold, color=text_color)
    add_textbox(s, Inches(0.6), Inches(5.7), Inches(12), Inches(1.4),
                "Source: WordStream travel benchmarks 2025 · SimilarWeb traffic data · Italy-localised CPC "
                "estimates from Google Keyword Planner. Our projected CPL of €6.67 is ~30% better than the "
                "closest competitor (Lastminute) — driven by tighter audience targeting (IT millennials only) "
                "and a higher-value lead magnet (€50 coupon vs. newsletter signup).",
                font_size=10, color=MUTED, italic=True, line_spacing=1.4)
    add_footer(s, PAGE)
slide_competitors()

# ---- ACCESSIBILITY
def slide_a11y():
    global PAGE; PAGE += 1
    s = new_content_slide("Accessibility audit", "WCAG 2.1 AA review",
                          "6 PASS · 3 PARTIAL · documented fixes · screen-reader friendly")
    a11y = [
        ("✓", "Color contrast", "Forest on cream = 12.4:1 (AAA) · Burnt on cream = 4.6:1 (AA)", "PASS"),
        ("✓", "Heading hierarchy", "Single h1 per page · h2/h3 cascade · semantic <nav>", "PASS"),
        ("✓", "Alt text", "All images alt'd · decorative use aria-hidden=true", "PASS"),
        ("◐", "Form labels", "All 5 fields labelled · could improve aria-required=true", "PARTIAL"),
        ("◐", "Keyboard nav", "Tab traversal works · sticky CTA may need skip-link", "PARTIAL"),
        ("◐", "Cookie banner", "Dismissible via Escape · should announce as aria-live", "PARTIAL"),
        ("✓", "Focus indicators", "Default focus rings preserved · visible everywhere", "PASS"),
        ("✓", "Lang attribute", "html lang='it' set · alternate hreflang for /?lang=en", "PASS"),
    ]
    y = 2.5
    for i, (icon, item, desc, status) in enumerate(a11y):
        yy = y + i * 0.5
        color = FOREST if status == "PASS" else BURNT
        add_textbox(s, Inches(0.6), Inches(yy), Inches(0.5), Inches(0.4),
                    icon, font_size=14, bold=True, color=color)
        add_textbox(s, Inches(1.1), Inches(yy + 0.05), Inches(2.5), Inches(0.4),
                    item, font_size=11, bold=True, color=FOREST)
        add_textbox(s, Inches(3.6), Inches(yy + 0.07), Inches(8), Inches(0.4),
                    desc, font_size=9, color=INK)
        add_textbox(s, Inches(11.7), Inches(yy + 0.08), Inches(1.0), Inches(0.4),
                    status, font_size=9, bold=True, color=color, align=PP_ALIGN.RIGHT)
    add_textbox(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
                "Tooling: Manual code review + axe-core DevTools. Fixes shipped to v2 backlog.",
                font_size=9, color=MUTED, italic=True)
    add_footer(s, PAGE)
slide_a11y()

# ===== SECTION 11 ===========================================
new_section_divider("11", "Deliverables",
                    "Vercel · GitHub · live URLs · what we shipped"); PAGE += 1

# ---- VERCEL DEPLOYMENTS
def slide_vercel():
    global PAGE; PAGE += 1
    s = new_content_slide("Live · Vercel deployments", "Auto-deploy on every push to main",
                          "12 deployments · all READY · sub-40 second build · public production URL")
    add_image_below_title(s, os.path.join(SHOTS, "vercel.png"),
                          top_in=2.4, max_w_in=11.5, max_h_in=4.4)
    add_footer(s, PAGE)
slide_vercel()

# ---- GITHUB
def slide_github():
    global PAGE; PAGE += 1
    s = new_content_slide("Live · GitHub repository", "Public · every commit auditable",
                          "Configuration as code: site, GTM JSON, briefings, Brevo templates · all in one repo")
    add_image_below_title(s, os.path.join(SHOTS, "github-repo.png"),
                          top_in=2.4, max_w_in=11.5, max_h_in=4.4)
    add_footer(s, PAGE)
slide_github()

# ---- WHAT WE DELIVERED
def slide_delivered():
    global PAGE; PAGE += 1
    s = new_content_slide("What we delivered", "Beyond the brief",
                          "Every brief requirement met · plus the infrastructure to launch tomorrow")
    items = [
        ("Live website", "Next.js + Vercel · backpack-bottle.vercel.app · public"),
        ("Brand identity system", "Wordmark · 6-color palette · Fraunces + Inter · voice guidelines"),
        ("5 destinations", "Rome · Barcelona · Amsterdam · Lisbon · Prague · individual landing pages"),
        ("BB50 coupon", "Branded PDF · €50 off · 90 days · validated"),
        ("Lead form + API", "Server-side validation · rate-limited · Brevo + Meta CAPI fan-out"),
        ("Brevo automation", "List 13 · 8-email welcome series · 7 personalisation attributes"),
        ("Tracking", "GTM container live · 17 tags · 22 triggers · GA4 receiving events"),
        ("Compliance", "Consent Mode v2 · cookie tracking · GDPR privacy + T&Cs"),
        ("3 SEO blog posts", "/blog · IT + EN · long-form · internal linking"),
        ("Measurement plan", "/measurement live · KPI tree · attribution model · audit-ready"),
        ("Stack page", "/stack · 6-layer architecture · live + browseable"),
        ("Looker Studio (mock)", "/dashboard · 6 panels · projection ready"),
    ]
    y = 2.5
    for i, (item, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.6 + col * 6.18
        yy = y + row * 0.65
        add_textbox(s, Inches(x), Inches(yy), Inches(0.3), Inches(0.3),
                    "✓", font_size=14, bold=True, color=BURNT)
        add_textbox(s, Inches(x + 0.4), Inches(yy + 0.02), Inches(2.5), Inches(0.3),
                    item, font_size=11, bold=True, color=FOREST)
        add_textbox(s, Inches(x + 2.9), Inches(yy + 0.05), Inches(3.2), Inches(0.3),
                    desc, font_size=9, color=MUTED)
    add_footer(s, PAGE)
slide_delivered()

# ---- WHERE TO FIND EVERYTHING
def slide_where():
    global PAGE; PAGE += 1
    s = new_content_slide("Where to find everything", "Live URLs · public · shareable",
                          "Bookmark these · they all work today")
    urls = [
        ("Live site", "backpack-bottle.vercel.app"),
        ("Hero / homepage", "backpack-bottle.vercel.app/"),
        ("Destination detail", "backpack-bottle.vercel.app/destinazioni/lisbon"),
        ("Thank-you page", "backpack-bottle.vercel.app/grazie"),
        ("Privacy & T&Cs", "backpack-bottle.vercel.app/privacy"),
        ("Blog (3 posts)", "backpack-bottle.vercel.app/blog"),
        ("Measurement plan", "backpack-bottle.vercel.app/measurement"),
        ("Stack architecture", "backpack-bottle.vercel.app/stack"),
        ("Looker mock", "backpack-bottle.vercel.app/dashboard"),
        ("GitHub repo", "github.com/saikanagat1117-sys/backpack-bottle"),
        ("Sitemap", "backpack-bottle.vercel.app/sitemap.xml"),
        ("Coupon PDF", "backpack-bottle.vercel.app/backpack-bottle-coupon-BB50.pdf"),
    ]
    y = 2.5
    for i, (label, url) in enumerate(urls):
        col = i % 2
        row = i // 2
        x = 0.6 + col * 6.18
        yy = y + row * 0.65
        add_rect(s, Inches(x), Inches(yy), Inches(5.95), Inches(0.55), WHITE,
                 line_color=CREAM_DARK)
        add_textbox(s, Inches(x + 0.2), Inches(yy + 0.08), Inches(2.0), Inches(0.4),
                    label, font_size=10, bold=True, color=BURNT)
        add_textbox(s, Inches(x + 0.2), Inches(yy + 0.3), Inches(5.7), Inches(0.3),
                    url, font_size=9, color=FOREST, font_name="Courier")
    add_footer(s, PAGE)
slide_where()

# ===== SECTION 12 ===========================================
new_section_divider("12", "Closing",
                    "Risks · what's next · thank you"); PAGE += 1

# ---- RISKS
def slide_risks():
    global PAGE; PAGE += 1
    s = new_content_slide("Risks & mitigations", "What could go wrong · what we'd do",
                          "Honest about uncertainty · plans on the shelf for each scenario")
    risks = [
        ("Meta CPMs spike",
         "Italian travel CPMs rose 22% in 2025 summer. If they spike again, we shift budget to retargeting (lower CPM) and reduce broad awareness."),
        ("Form conversion below target",
         "If form_submit < 5% by week 2, drop the 'departure city' field (cuts to 4 fields). A/B variant B (emotional CTA) moves to 100% rollout."),
        ("Brevo deliverability",
         "Free tier with new sender domain. If open rate <20%, we warm-up with 100 contacts before the campaign opens, then ramp."),
        ("iOS 14+ attribution loss",
         "Already mitigated: Meta CAPI server-side mirror covers ~30% of attribution that client-side Pixel misses."),
        ("Coupon abuse",
         "BB50 has min basket €300 + one-per-customer enforcement (Brevo dedupe by email + IP). Risk is low but monitored daily."),
        ("Google Ads policy review",
         "Travel + competitor bidding can trigger reviews. Pre-flight ad copy through policy checker, hold backup ad sets."),
    ]
    y = 2.5
    for i, (risk, mit) in enumerate(risks):
        col = i % 2
        row = i // 2
        x = 0.6 + col * 6.18
        yy = y + row * 1.55
        add_rect(s, Inches(x), Inches(yy), Inches(5.95), Inches(1.4), WHITE,
                 line_color=CREAM_DARK)
        add_rect(s, Inches(x), Inches(yy), Inches(0.15), Inches(1.4), BURNT)
        add_textbox(s, Inches(x + 0.3), Inches(yy + 0.15), Inches(5.5), Inches(0.3),
                    risk, font_size=11, bold=True, color=FOREST)
        add_textbox(s, Inches(x + 0.3), Inches(yy + 0.5), Inches(5.5), Inches(0.85),
                    mit, font_size=9.5, color=INK, line_spacing=1.3)
    add_footer(s, PAGE)
slide_risks()

# ---- THANK YOU
def slide_thanks():
    global PAGE; PAGE += 1
    s = prs.slides.add_slide(BLANK)
    add_filled_bg(s, FOREST)
    add_textbox(s, Inches(0.6), Inches(0.5), Inches(11), Inches(0.3),
                "GROUP 3 · DIGITAL PLATFORMS LAB · BOLOGNA BUSINESS SCHOOL",
                font_size=9, bold=True, color=BURNT)
    add_textbox(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5),
                "Thank you.",
                font_size=84, color=CREAM, font_name="Times-Bold")
    add_textbox(s, Inches(1), Inches(3.7), Inches(11.3), Inches(0.6),
                "Skip the research fatigue.",
                font_size=22, color=BURNT, italic=True, font_name="Times-Bold")

    add_textbox(s, Inches(1), Inches(5.0), Inches(2.5), Inches(0.3),
                "GROUP 3", font_size=10, bold=True, color=BURNT)
    add_textbox(s, Inches(1), Inches(5.35), Inches(11), Inches(0.4),
                "Sai Prathyaksh Kanagat  ·  Marica Motta  ·  Alberto Faggiotto",
                font_size=14, color=CREAM)
    add_textbox(s, Inches(1), Inches(5.75), Inches(11), Inches(0.4),
                "Cecilia Gullett  ·  Marilu Stevens  ·  Amelia Ruben",
                font_size=14, color=CREAM)

    add_textbox(s, Inches(1), Inches(6.55), Inches(11), Inches(0.3),
                "Live: backpack-bottle.vercel.app  ·  Code: github.com/saikanagat1117-sys/backpack-bottle",
                font_size=10, color=BURNT)
slide_thanks()

# =================================================================
# Replace {TOTAL} placeholder in all footers
# =================================================================
total_slides = len(prs.slides)
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if PAGE_TOTAL_PLACEHOLDER in run.text:
                    run.text = run.text.replace(PAGE_TOTAL_PLACEHOLDER, str(total_slides))

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {total_slides}")
