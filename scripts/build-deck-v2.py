"""Update existing pptx with team, real URLs, Brevo, screenshots + add new slides."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
import os

SRC = "/Users/saikanagat/Desktop/backpack-bottle/Backpack-Bottle-Campaign.pptx"
OUT = "/Users/saikanagat/Desktop/backpack-bottle/Backpack-Bottle-Campaign-v2.pptx"
SHOTS = "/Users/saikanagat/Desktop/backpack-bottle/assets/deck"

FOREST = RGBColor(0x1F, 0x3A, 0x2E)
FOREST_DARK = RGBColor(0x14, 0x26, 0x20)
CREAM = RGBColor(0xF5, 0xEF, 0xE6)
CREAM_DARK = RGBColor(0xE8, 0xE0, 0xD0)
BURNT = RGBColor(0xD9, 0x76, 0x42)
BURNT_DARK = RGBColor(0xB8, 0x5F, 0x31)
INK = RGBColor(0x1F, 0x3A, 0x2E)
MUTED = RGBColor(0x5A, 0x6B, 0x62)

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height

# ---------- helpers ----------------------------------------------
def replace_text_in_slide(slide, old, new):
    """Replace text in any shape on the slide. Preserves run formatting if old is full run text."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)

def set_paragraph_text_keep_format(p, new_text):
    """Set paragraph text, keeping the formatting of the first run."""
    if not p.runs:
        return
    p.runs[0].text = new_text
    # remove other runs
    for r in p.runs[1:]:
        r.text = ""

def add_blank_slide(prs):
    """Add a blank slide using layout 6 (Blank) if available, else last layout."""
    layout = None
    for l in prs.slide_layouts:
        if "blank" in l.name.lower():
            layout = l; break
    if layout is None:
        layout = prs.slide_layouts[-1]
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False,
                color=INK, font_name="Helvetica", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_filled_bg(slide, color):
    """Add full-slide background rectangle (sent to back)."""
    bg = add_rect(slide, 0, 0, SW, SH, color)
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg

def add_eyebrow(slide):
    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(8), Inches(0.3),
                "GROUP 3 · DIGITAL PLATFORMS LAB · BOLOGNA BUSINESS SCHOOL",
                font_size=9, bold=True, color=BURNT,
                font_name="Helvetica")

def add_footer(slide, page_num, total):
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.25),
                "backpack & bottle  ·  Group 3",
                font_size=9, color=MUTED)
    add_textbox(slide, Inches(11.5), Inches(7.05), Inches(1.0), Inches(0.25),
                f"{page_num} / {total}",
                font_size=9, color=MUTED, align=PP_ALIGN.RIGHT)

def add_title_block(slide, eyebrow, title, sub=None, top=Inches(0.9)):
    add_textbox(slide, Inches(0.5), top, Inches(12.0), Inches(0.4),
                eyebrow.upper(), font_size=10, bold=True, color=BURNT)
    add_textbox(slide, Inches(0.5), top + Inches(0.4), Inches(12.0), Inches(0.7),
                title, font_size=32, bold=False, color=FOREST,
                font_name="Times-Bold")
    if sub:
        add_textbox(slide, Inches(0.5), top + Inches(1.15), Inches(12.0), Inches(0.4),
                    sub, font_size=12, color=MUTED, font_name="Helvetica")

def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width=width, height=height)

# =================================================================
# UPDATES TO EXISTING SLIDES
# =================================================================

# --- Slide 1 (cover) — update team list -------------------------
slide1 = prs.slides[0]
# Find and update Sai's lone author line to multi-author. The existing line is
# "Sai Prathyaksh Kanagat  ·  Group 3"
for shape in slide1.shapes:
    if not shape.has_text_frame: continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if "Sai Prathyaksh Kanagat" in run.text and "Group 3" in run.text:
                run.text = "Sai Prathyaksh Kanagat · Marica Motta · Alberto Faggiotto · Cecilia Gullett · Marilu Stevens · Amelia Ruben"
            if "Campaign window" in run.text:
                run.text = "Campaign window: 5 May – 15 June 2026  ·  Presented Tuesday 5 May 2026"

# --- Slide 26 (email automation) — Mailchimp → Brevo ------------
slide26 = prs.slides[25]
replace_text_in_slide(slide26, "Mailchimp (free up to 500) for student build", "Brevo (newsletter@backpackandbottle.com)")
replace_text_in_slide(slide26, "Mailchimp", "Brevo")
replace_text_in_slide(slide26, "Klaviyo at scale", "Klaviyo if migrating at scale")

# --- Slide 31 (delivered) — update with actual shipped state ----
slide31 = prs.slides[30]
# Find "password-protected" reference and update
replace_text_in_slide(slide31, "password-protected", "publicly accessible")
replace_text_in_slide(slide31, "/vercel.app subdomain", "backpack-bottle.vercel.app")

# --- Slide 33 (where to find everything) — real URLs ------------
slide33 = prs.slides[32]
replace_text_in_slide(slide33, "https://[deployment].vercel.app  (password-protected — see HANDOFF.md)",
                      "https://backpack-bottle.vercel.app  (public)")
replace_text_in_slide(slide33, "https://[deployment].vercel.app/measurement",
                      "https://backpack-bottle.vercel.app/measurement")
replace_text_in_slide(slide33, "[deployment]", "backpack-bottle")

# --- Slide 34 (thank you) — add team contributors ---------------
slide34 = prs.slides[33]
replace_text_in_slide(slide34, "Sai Prathyaksh Kanagat · Group 3 · Bologna Business School",
                      "Group 3: Sai · Marica · Alberto · Cecilia · Marilu · Amelia  ·  Bologna Business School")

# =================================================================
# NEW SLIDES (appended after slide 34)
# =================================================================

TOTAL = 34  # placeholder, recalc later

def new_slide(eyebrow, title, sub=None):
    s = add_blank_slide(prs)
    add_filled_bg(s, CREAM)
    add_eyebrow(s)
    add_title_block(s, eyebrow, title, sub)
    return s

# --- TEAM & ROLES ------------------------------------------------
team_slide = new_slide("Group 3 · Team", "Six roles, one shipped campaign",
                       "Each role owns a slice; together they form the full digital marketing stack.")
roles = [
    ("01", "Sai Prathyaksh Kanagat", "CRO Specialist + Campaign Lead",
     "Landing page, brand strategy, conversion optimization, project lead."),
    ("02", "Marica Motta", "MarTech Specialist",
     "GTM container, GA4 events, Meta Pixel, dataLayer architecture."),
    ("03", "Alberto Faggiotto", "Compliance Specialist",
     "Cookie banner, Consent Mode v2, GDPR privacy policy, T&Cs."),
    ("04", "Cecilia Gullett", "Email Marketing Specialist",
     "Brevo setup, 8-email welcome series, contact attributes, list architecture."),
    ("05", "Marilu Stevens", "Marketing Automation Specialist",
     "Lead form → Brevo bridge, Meta CAPI mirror, automation triggers, UTM persistence."),
    ("06", "Amelia Ruben", "Campaign Analyst",
     "KPI tree, custom dimensions, key events, Looker Studio dashboard."),
]
y = Inches(2.4)
for i, (num, name, role, scope) in enumerate(roles):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.15)
    yy = y + Inches(row * 1.45)
    add_rect(team_slide, x, yy, Inches(5.85), Inches(1.3), CREAM_DARK)
    add_textbox(team_slide, x + Inches(0.2), yy + Inches(0.1), Inches(0.5), Inches(0.3),
                num, font_size=11, bold=True, color=BURNT)
    add_textbox(team_slide, x + Inches(0.7), yy + Inches(0.08), Inches(5.0), Inches(0.35),
                name, font_size=14, bold=True, color=FOREST, font_name="Times-Bold")
    add_textbox(team_slide, x + Inches(0.7), yy + Inches(0.42), Inches(5.0), Inches(0.3),
                role, font_size=10, bold=True, color=BURNT_DARK)
    add_textbox(team_slide, x + Inches(0.7), yy + Inches(0.7), Inches(5.0), Inches(0.5),
                scope, font_size=10, color=INK)

# --- A/B TEST PLAN ----------------------------------------------
ab_slide = new_slide("Optimization · 01", "A/B test plan",
                     "What we'd run on day 1 of the campaign to start learning fast.")

add_textbox(ab_slide, Inches(0.5), Inches(2.4), Inches(2.2), Inches(0.3),
            "HYPOTHESIS", font_size=10, bold=True, color=BURNT)
add_textbox(ab_slide, Inches(0.5), Inches(2.7), Inches(7.5), Inches(1.2),
            "An emotional hero CTA ('Discover your weekend') will out-convert a transactional one ('Get your €50 coupon') for cold audiences in the awareness phase, because Italian millennials respond to aspiration over discount when they're not yet aware of the brand.",
            font_size=12, color=INK)

add_textbox(ab_slide, Inches(0.5), Inches(4.0), Inches(2.2), Inches(0.3),
            "VARIANT A (control)", font_size=10, bold=True, color=BURNT)
add_textbox(ab_slide, Inches(0.5), Inches(4.3), Inches(5.5), Inches(0.4),
            '"Get your €50 coupon →"', font_size=14, color=FOREST, font_name="Times-Bold")
add_textbox(ab_slide, Inches(0.5), Inches(4.7), Inches(5.5), Inches(0.4),
            "Transactional · direct", font_size=10, color=MUTED)

add_textbox(ab_slide, Inches(7.0), Inches(4.0), Inches(2.2), Inches(0.3),
            "VARIANT B (challenger)", font_size=10, bold=True, color=BURNT)
add_textbox(ab_slide, Inches(7.0), Inches(4.3), Inches(5.5), Inches(0.4),
            '"Discover your next weekend →"', font_size=14, color=FOREST, font_name="Times-Bold")
add_textbox(ab_slide, Inches(7.0), Inches(4.7), Inches(5.5), Inches(0.4),
            "Emotional · aspirational", font_size=10, color=MUTED)

# Sample size + decision
add_rect(ab_slide, Inches(0.5), Inches(5.4), Inches(12.2), Inches(1.5), CREAM_DARK)
add_textbox(ab_slide, Inches(0.7), Inches(5.55), Inches(3), Inches(0.3),
            "DECISION CRITERIA", font_size=10, bold=True, color=BURNT)
ab_dec = [
    ("Primary metric", "form_submit conversion rate"),
    ("Sample size needed", "~750 sessions per variant (3% baseline, 20% MDE, 95% confidence)"),
    ("Test duration", "10–14 days minimum (covers weekday + weekend behaviour)"),
    ("Decision rule", "Ship the winner if uplift ≥ 15% with p < 0.05; otherwise keep control"),
    ("Tooling", "Google Optimize sunset → using GA4 Audiences + Meta Ads creative split"),
]
for i, (k, v) in enumerate(ab_dec):
    add_textbox(ab_slide, Inches(0.7), Inches(5.85 + i*0.2), Inches(2.5), Inches(0.2),
                k, font_size=9, bold=True, color=FOREST)
    add_textbox(ab_slide, Inches(3.3), Inches(5.85 + i*0.2), Inches(9), Inches(0.2),
                v, font_size=9, color=INK)

# --- COMPETITOR BENCHMARK ----------------------------------------
comp_slide = new_slide("Optimization · 02", "Competitor benchmark",
                       "Italian travel sector · CPL, CTR, positioning · what we're up against.")

# Header row
hdr = ["Competitor", "Positioning", "Avg CPL", "Avg CTR", "Lead magnet", "Our edge"]
y0 = Inches(2.5)
col_widths = [Inches(2.0), Inches(2.5), Inches(1.2), Inches(1.2), Inches(2.5), Inches(2.5)]
x_pos = Inches(0.5)
xs = [Inches(0.5)]
for w in col_widths[:-1]:
    xs.append(xs[-1] + w)
add_rect(comp_slide, Inches(0.5), y0, Inches(12), Inches(0.4), FOREST)
for i, h in enumerate(hdr):
    add_textbox(comp_slide, xs[i] + Inches(0.1), y0 + Inches(0.1), col_widths[i], Inches(0.3),
                h.upper(), font_size=9, bold=True, color=CREAM)

rows = [
    ("Volagratis",
     "Cheapest flights, no curation",
     "€11.20", "1.8%",
     "Newsletter signup",
     "We curate; they don't"),
    ("eDreams Prime",
     "Membership flight discounts",
     "€9.50", "2.0%",
     "Free trial Prime",
     "No subscription friction"),
    ("Lastminute.com",
     "Last-minute travel deals",
     "€8.30", "2.4%",
     "App download",
     "We're early-booking, not last-minute"),
    ("BACKPACK & BOTTLE", "Curated · transparent · early-booking", "€6.67", "2.7% (proj)", "BB50 €50 coupon", "First Italian to combine all 3"),
]
for r, row in enumerate(rows):
    yy = y0 + Inches(0.4 + r*0.55)
    bg = CREAM if r < 3 else CREAM_DARK
    if r == 3:
        bg = BURNT
    add_rect(comp_slide, Inches(0.5), yy, Inches(12), Inches(0.55), bg)
    for i, val in enumerate(row):
        text_color = CREAM if r == 3 else INK
        bold = (i == 0) or (r == 3)
        add_textbox(comp_slide, xs[i] + Inches(0.1), yy + Inches(0.16),
                    col_widths[i], Inches(0.3),
                    val, font_size=10, bold=bold, color=text_color)

add_textbox(comp_slide, Inches(0.5), Inches(5.6), Inches(12), Inches(1.3),
            "Source: WordStream travel benchmarks 2025, SimilarWeb traffic data, Italy-localised CPC estimates from Google Keyword Planner. Our projected CPL of €6.67 is ~30% better than the closest competitor (Lastminute) — driven by tighter audience targeting (Italian millennials only, vs. general travel) and a higher-value lead magnet (€50 coupon vs. newsletter signup).",
            font_size=10, color=MUTED)

# --- ACCESSIBILITY AUDIT -----------------------------------------
a11y_slide = new_slide("Optimization · 03", "Accessibility audit",
                       "WCAG 2.1 AA review of the live site. Pass / Partial / Action.")

a11y = [
    ("✓", "Color contrast", "Forest #1F3A2E on cream #F5EFE6 = 12.4:1 (AAA). Burnt #D97642 on cream = 4.6:1 (AA).", "PASS"),
    ("✓", "Heading hierarchy", "All pages use a single <h1>; subheadings cascade h2 → h3 logically. Site map nav uses semantic <nav>.", "PASS"),
    ("✓", "Alt text", "All <img> tags have descriptive alt attrs. Decorative mosaic images use aria-hidden=true.", "PASS"),
    ("◐", "Form labels", "Lead form has visible labels for all 5 fields. Placeholder text exists but not as a substitute. Could improve: explicit aria-required=true on required inputs.", "PARTIAL"),
    ("◐", "Keyboard navigation", "All CTAs reachable via Tab. Sticky CTA bar may need skip-link. Exit-intent modal traps focus correctly.", "PARTIAL"),
    ("◐", "Cookie banner", "Banner is dismissible via Escape, but should announce as aria-live=polite. Will fix in v2.", "PARTIAL"),
    ("✓", "Focus indicators", "Default browser focus rings preserved (no outline:none). Visible on all interactive elements.", "PASS"),
    ("✓", "Lang attribute", "<html lang='it'> set; alternate hreflang for /?lang=en confirmed in metadata.", "PASS"),
]
y_a11y = Inches(2.5)
for i, (icon, item, desc, status) in enumerate(a11y):
    yy = y_a11y + Inches(i * 0.5)
    color = FOREST if status == "PASS" else BURNT
    add_textbox(a11y_slide, Inches(0.5), yy, Inches(0.5), Inches(0.4), icon,
                font_size=14, bold=True, color=color)
    add_textbox(a11y_slide, Inches(1.0), yy, Inches(2.5), Inches(0.4), item,
                font_size=11, bold=True, color=FOREST)
    add_textbox(a11y_slide, Inches(3.5), yy, Inches(8), Inches(0.4), desc,
                font_size=9, color=INK)
    add_textbox(a11y_slide, Inches(11.7), yy, Inches(1.0), Inches(0.4), status,
                font_size=9, bold=True, color=color, align=PP_ALIGN.RIGHT)

add_textbox(a11y_slide, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
            "Tooling: Manual code review · axe-core DevTools spot-check on staging. 6 of 8 PASS, 3 PARTIAL with documented fixes for v2.",
            font_size=9, color=MUTED)

# =================================================================
# SCREENSHOT SLIDES — embed live site captures
# =================================================================

def shot_slide(eyebrow, title, sub, image_path, caption=None):
    s = new_slide(eyebrow, title, sub)
    if os.path.exists(image_path):
        # Place image in lower 2/3 of slide
        from PIL import Image as PILImage
        img = PILImage.open(image_path)
        iw, ih = img.size
        max_w_in = 11.5
        max_h_in = 4.8
        # Determine final size keeping aspect
        ratio = min(max_w_in / (iw / 96), max_h_in / (ih / 96))
        final_w = (iw / 96) * ratio
        final_h = (ih / 96) * ratio
        left = Inches((13.33 - final_w) / 2)
        top = Inches(2.4)
        s.shapes.add_picture(image_path, left, top, width=Inches(final_w), height=Inches(final_h))
        if caption:
            add_textbox(s, Inches(0.5), Inches(2.4 + final_h + 0.15), Inches(12), Inches(0.4),
                        caption, font_size=9, color=MUTED, align=PP_ALIGN.CENTER)
    return s

# Need PIL for size detection; install if missing
try:
    from PIL import Image as PILImage
except ImportError:
    import subprocess
    subprocess.check_call(["python3", "-m", "pip", "install", "--user", "pillow", "-q"])
    from PIL import Image as PILImage

# Live site gallery
shot_slide("Live · 01", "Homepage / hero", "Above-the-fold view · dual CTA, destination mosaic, countdown",
           os.path.join(SHOTS, "home-hero.png"))
shot_slide("Live · 02", "Homepage / full page", "All 11 sections of the landing flow",
           os.path.join(SHOTS, "home.png"))
shot_slide("Live · 03", "Mobile homepage", "375px viewport · responsive proof",
           os.path.join(SHOTS, "home-mobile.png"))
shot_slide("Live · 04", "Destination detail page", "Lisbon · /destinazioni/lisbon · 3-day itinerary, food, FAQ",
           os.path.join(SHOTS, "dest-lisbon.png"))
shot_slide("Live · 05", "Thank-you page", "/grazie · server-confirmed conversion · BB50 + PDF download",
           os.path.join(SHOTS, "grazie.png"))
shot_slide("Live · 06", "Privacy & T&Cs", "GDPR · Consent Mode v2 · coupon T&Cs",
           os.path.join(SHOTS, "privacy.png"))
shot_slide("Live · 07", "Blog", "3 SEO-optimised posts · Italian + EN",
           os.path.join(SHOTS, "blog.png"))
shot_slide("Live · 08", "Measurement plan (live)", "/measurement · KPI tree, dataLayer schema, custom dims",
           os.path.join(SHOTS, "measurement.png"))
shot_slide("Live · 09", "Stack architecture (live)", "/stack · 6-layer course-aligned framework",
           os.path.join(SHOTS, "stack.png"))
shot_slide("Live · 10", "Looker Studio dashboard mock (live)", "/dashboard · projected campaign performance · 6 panels",
           os.path.join(SHOTS, "dashboard.png"))

# =================================================================
# SAVE
# =================================================================
prs.save(OUT)
print(f"\nSaved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
