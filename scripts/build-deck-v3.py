"""Deck v3 — clean flow, screenshots inline, stack visual, per-role depth."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
from PIL import Image as PILImage

SRC = "/Users/saikanagat/Desktop/backpack-bottle/Backpack-Bottle-Campaign.pptx"
OUT = "/Users/saikanagat/Desktop/backpack-bottle/Backpack-Bottle-Campaign-v3.pptx"
SHOTS = "/Users/saikanagat/Desktop/backpack-bottle/assets/deck"

FOREST = RGBColor(0x1F, 0x3A, 0x2E)
FOREST_DARK = RGBColor(0x14, 0x26, 0x20)
CREAM = RGBColor(0xF5, 0xEF, 0xE6)
CREAM_DARK = RGBColor(0xE8, 0xE0, 0xD0)
BURNT = RGBColor(0xD9, 0x76, 0x42)
BURNT_DARK = RGBColor(0xB8, 0x5F, 0x31)
INK = RGBColor(0x1F, 0x3A, 0x2E)
MUTED = RGBColor(0x5A, 0x6B, 0x62)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Open existing as base for visual style continuity
prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height

# Strategy:
# 1. Edit specific slides in-place (text replacements + image insertions)
# 2. Replace whacky Looker slide content with placeholder
# 3. INSERT new slides at strategic points using copy of existing layout
# 4. Add stack visual as NEW slide right after slide 14 (site map)

# ---------- helpers ----------------------------------------------
def replace_text(slide, old, new):
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)

def add_blank_slide():
    layout = next((l for l in prs.slide_layouts if "blank" in l.name.lower()), prs.slide_layouts[-1])
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False,
                color=INK, font_name="Helvetica", align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width: shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_filled_bg(slide, color):
    bg = add_rect(slide, 0, 0, SW, SH, color)
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg

def add_eyebrow(slide):
    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(8), Inches(0.3),
                "GROUP 3 · DIGITAL PLATFORMS LAB · BOLOGNA BUSINESS SCHOOL",
                font_size=9, bold=True, color=BURNT)

def add_footer(slide, page_num):
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.25),
                "backpack & bottle  ·  Group 3",
                font_size=9, color=MUTED)

def add_title_block(slide, eyebrow, title, sub=None, top=Inches(0.85)):
    add_textbox(slide, Inches(0.5), top, Inches(12.0), Inches(0.4),
                eyebrow.upper(), font_size=10, bold=True, color=BURNT)
    add_textbox(slide, Inches(0.5), top + Inches(0.4), Inches(12.0), Inches(0.7),
                title, font_size=30, bold=False, color=FOREST,
                font_name="Times-Bold")
    if sub:
        add_textbox(slide, Inches(0.5), top + Inches(1.1), Inches(12.0), Inches(0.4),
                    sub, font_size=12, color=MUTED)

def section_slide(eyebrow, title, sub=None):
    s = add_blank_slide()
    add_filled_bg(s, CREAM)
    add_eyebrow(s)
    add_title_block(s, eyebrow, title, sub)
    return s

def insert_image_to_slide(slide, image_path, left_in, top_in, max_w_in=None, max_h_in=None,
                          shadow=True):
    """Insert image keeping aspect ratio, optional shadow."""
    if not os.path.exists(image_path): return None
    img = PILImage.open(image_path)
    iw, ih = img.size
    aspect = ih / iw
    if max_w_in and max_h_in:
        if iw / 96 * (max_h_in / (ih / 96)) <= max_w_in:
            final_h = max_h_in
            final_w = max_h_in / aspect
        else:
            final_w = max_w_in
            final_h = max_w_in * aspect
    elif max_w_in:
        final_w = max_w_in
        final_h = max_w_in * aspect
    elif max_h_in:
        final_h = max_h_in
        final_w = max_h_in / aspect
    else:
        final_w, final_h = 6, 6 * aspect
    pic = slide.shapes.add_picture(image_path, Inches(left_in), Inches(top_in),
                                    width=Inches(final_w), height=Inches(final_h))
    return pic

# Helper: fill an existing slide with a screenshot below the title
def add_screenshot_to_existing_slide(slide_idx, image_path, top_in=2.4, max_w_in=11.5, max_h_in=4.5):
    """Add screenshot to existing slide (centred horizontally below title)."""
    if not os.path.exists(image_path): return
    slide = prs.slides[slide_idx]
    img = PILImage.open(image_path)
    aspect = img.size[1] / img.size[0]
    final_w = min(max_w_in, max_h_in / aspect)
    final_h = final_w * aspect
    if final_h > max_h_in:
        final_h = max_h_in
        final_w = final_h / aspect
    left = (13.33 - final_w) / 2
    slide.shapes.add_picture(image_path, Inches(left), Inches(top_in),
                             width=Inches(final_w), height=Inches(final_h))

# =================================================================
# 1. UPDATE EXISTING SLIDE TEXT (same as v2)
# =================================================================
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if not shape.has_text_frame: continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if "Sai Prathyaksh Kanagat" in run.text and "Group 3" in run.text:
                run.text = "Group 3 · Sai Prathyaksh Kanagat · Marica Motta · Alberto Faggiotto · Cecilia Gullett · Marilu Stevens · Amelia Ruben"
            if "Campaign window" in run.text:
                run.text = "Campaign window: 5 May – 15 June 2026  ·  Presented Tuesday 5 May 2026"

slide26 = prs.slides[25]
replace_text(slide26, "Mailchimp (free up to 500) for student build", "Brevo (newsletter@backpackandbottle.com)")
replace_text(slide26, "Mailchimp", "Brevo")
replace_text(slide26, "Klaviyo at scale", "Klaviyo if migrating at scale")

slide31 = prs.slides[30]
replace_text(slide31, "password-protected", "publicly accessible")
replace_text(slide31, "/vercel.app subdomain", "backpack-bottle.vercel.app")

slide33 = prs.slides[32]
replace_text(slide33, "https://[deployment].vercel.app  (password-protected — see HANDOFF.md)",
             "https://backpack-bottle.vercel.app  (public)")
replace_text(slide33, "https://[deployment].vercel.app/measurement",
             "https://backpack-bottle.vercel.app/measurement")
replace_text(slide33, "[deployment]", "backpack-bottle")

slide34 = prs.slides[33]
replace_text(slide34, "Sai Prathyaksh Kanagat · Group 3 · Bologna Business School",
             "Group 3: Sai · Marica · Alberto · Cecilia · Marilu · Amelia  ·  Bologna Business School")

# =================================================================
# 2. INSERT SCREENSHOTS ONTO EXISTING SLIDES (where there's space)
# =================================================================
# Slide 10 — The offer (5 cities): add destinations-grid screenshot
add_screenshot_to_existing_slide(9, os.path.join(SHOTS, "destinations-grid.png"),
                                  top_in=4.5, max_w_in=12, max_h_in=2.5)

# Slide 13 — Landing page architecture: add home-hero
add_screenshot_to_existing_slide(12, os.path.join(SHOTS, "home-hero.png"),
                                  top_in=4.0, max_w_in=11.5, max_h_in=3.0)

# Slide 14 — Site map: add github-repo screenshot
add_screenshot_to_existing_slide(13, os.path.join(SHOTS, "github-repo.png"),
                                  top_in=2.4, max_w_in=11.5, max_h_in=4.5)

# Slide 18 — Consent Mode v2: I'll handle below by replacing it with new content

# Slide 27 — Looker Studio dashboard: REPLACE with placeholder
# Easiest: clear all shapes and rebuild
slide27 = prs.slides[26]
# Remove existing text shapes (keep first textbox = title block)
shapes_to_remove = []
for shape in list(slide27.shapes):
    shapes_to_remove.append(shape)
for shape in shapes_to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

add_filled_bg(slide27, CREAM)
add_eyebrow(slide27)
add_title_block(slide27, "Reporting", "Looker Studio dashboard",
                "What we look at every morning during the campaign")
# Placeholder block
ph = add_rect(slide27, Inches(2), Inches(2.6), Inches(9.33), Inches(4.0),
              CREAM_DARK, line_color=BURNT)
add_textbox(slide27, Inches(2), Inches(3.3), Inches(9.33), Inches(0.5),
            "[ Live dashboard screenshot to be inserted on launch day ]",
            font_size=14, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
add_textbox(slide27, Inches(2), Inches(4.0), Inches(9.33), Inches(0.5),
            "Built in Looker Studio · 6 panels · refreshes every 12 hours",
            font_size=11, color=FOREST, align=PP_ALIGN.CENTER)
add_textbox(slide27, Inches(2), Inches(4.5), Inches(9.33), Inches(2),
            "Acquisition  ·  Funnel  ·  Cost  ·  Audience  ·  Engagement  ·  Coupon",
            font_size=12, bold=True, color=BURNT, align=PP_ALIGN.CENTER)
add_textbox(slide27, Inches(2), Inches(5.3), Inches(9.33), Inches(1.0),
            "Connected sources: GA4 (G-FHBQE8QZM9), Google Ads, Meta Ads via Supermetrics, "
            "Brevo for offline conversions. Shareable view-only link will be added once campaign launches.",
            font_size=10, color=INK, align=PP_ALIGN.CENTER)

# =================================================================
# 3. NEW SLIDES — strategically positioned by appending then reordering NOT done.
#    For simplicity: append to end. Existing flow already works; we treat new
#    slides as appendices grouped by theme.
# =================================================================

# ============= SECTION A: STACK ARCHITECTURE VISUAL ============
def build_stack_diagram_slide():
    s = section_slide("Architecture", "The full stack",
                      "Six layers · 27 components · single source of truth")

    # Layout: 6 horizontal bands, each labeled, with tool pills inside.
    layers = [
        ("01 · ACQUISITION", "How users arrive", FOREST,
         ["Meta Ads (60%)", "Google Ads (40%)", "Organic SEO", "Newsletter strip"]),
        ("02 · ENGAGEMENT", "What they see", FOREST_DARK,
         ["Next.js 14 / Vercel", "IT/EN locale", "Microsoft Clarity", "Exit-intent modal", "WhatsApp button"]),
        ("03 · CONVERSION", "How they act", BURNT,
         ["Lead form (5 fields)", "BB50 coupon · €50", "Countdown timer", "Sticky CTA", "/grazie"]),
        ("04 · RETENTION", "How they come back", BURNT_DARK,
         ["Brevo ESP", "8-email automation #23", "UTM persistence", "Meta CAPI mirror"]),
        ("05 · MEASUREMENT", "How we know it worked", FOREST,
         ["GTM-KSN24TZ2 (v5)", "GA4: G-FHBQE8QZM9", "Meta Pixel", "Looker Studio", "Consent Mode v2"]),
        ("06 · INFRA", "Where it lives", MUTED,
         ["Vercel (auto-deploy)", "GitHub (public)", "Next.js metadata API", "Rate limiting + GDPR"]),
    ]

    y0 = 2.2
    band_h = 0.72
    total_w = 12.33
    label_w = 2.2
    for i, (label, sub, accent, tools) in enumerate(layers):
        y = y0 + i * (band_h + 0.05)
        # Band background
        add_rect(s, Inches(0.5), Inches(y), Inches(total_w), Inches(band_h),
                 CREAM_DARK, line_color=accent)
        # Left label panel
        add_rect(s, Inches(0.5), Inches(y), Inches(label_w), Inches(band_h), accent)
        add_textbox(s, Inches(0.6), Inches(y + 0.08), Inches(label_w - 0.2), Inches(0.3),
                    label, font_size=10, bold=True, color=CREAM)
        add_textbox(s, Inches(0.6), Inches(y + 0.32), Inches(label_w - 0.2), Inches(0.3),
                    sub, font_size=8, color=CREAM, italic=True)
        # Tools pills
        x = 0.5 + label_w + 0.15
        pill_h = 0.32
        pill_y = y + (band_h - pill_h) / 2
        for tool in tools:
            tw = max(1.4, len(tool) * 0.085)
            if x + tw > 0.5 + total_w - 0.1:
                break
            add_rect(s, Inches(x), Inches(pill_y), Inches(tw), Inches(pill_h),
                     WHITE, line_color=accent)
            add_textbox(s, Inches(x + 0.08), Inches(pill_y + 0.05), Inches(tw - 0.15), Inches(0.25),
                        tool, font_size=8.5, color=FOREST, align=PP_ALIGN.CENTER)
            x += tw + 0.1

    # Caption strip
    add_textbox(s, Inches(0.5), Inches(6.65), Inches(12), Inches(0.4),
                "User journey ▸  ad → site → form → coupon → email → booking · all measured · all consent-gated",
                font_size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    return s

build_stack_diagram_slide()

# ============= SECTION B: 5 CITIES — INDIVIDUAL CARDS ===========
def build_cities_grid_slide():
    s = section_slide("Live · The offer", "5 curated city breaks",
                      "1.5–3 hour flights from Italy · transparent all-in pricing")
    cities = [
        ("Rome", "ROM", "rome", "€280", "2n"),
        ("Barcelona", "BCN", "barcelona", "€320", "3n"),
        ("Amsterdam", "AMS", "amsterdam", "€365", "3n"),
        ("Lisbon", "LIS", "lisbon", "€345", "3n"),
        ("Prague", "PRG", "prague", "€295", "2n"),
    ]
    cw = 2.36
    gap = 0.12
    total_grid_w = 5 * cw + 4 * gap
    x0 = (13.33 - total_grid_w) / 2
    y = 2.5
    for i, (name, code, slug, price, nights) in enumerate(cities):
        cx = x0 + i * (cw + gap)
        # Card bg
        add_rect(s, Inches(cx), Inches(y), Inches(cw), Inches(3.6),
                 WHITE, line_color=CREAM_DARK)
        # Image — try to insert
        img_path = os.path.join(SHOTS, f"dest-{slug}.png")
        if os.path.exists(img_path):
            # Fit a thumbnail at top of card
            img = PILImage.open(img_path)
            iw, ih = img.size
            target_w = cw - 0.2
            target_h = 1.6
            # Crop center horizontal slice
            s.shapes.add_picture(img_path, Inches(cx + 0.1), Inches(y + 0.1),
                                  width=Inches(target_w), height=Inches(target_h))
        # Card content
        add_textbox(s, Inches(cx + 0.15), Inches(y + 1.85), Inches(cw - 0.3), Inches(0.3),
                    code, font_size=9, bold=True, color=BURNT)
        add_textbox(s, Inches(cx + 0.15), Inches(y + 2.1), Inches(cw - 0.3), Inches(0.5),
                    name, font_size=18, bold=True, color=FOREST, font_name="Times-Bold")
        add_textbox(s, Inches(cx + 0.15), Inches(y + 2.7), Inches(cw - 0.3), Inches(0.3),
                    f"From {price} · {nights}", font_size=11, bold=True, color=BURNT_DARK)
        add_textbox(s, Inches(cx + 0.15), Inches(y + 3.1), Inches(cw - 0.3), Inches(0.4),
                    f"backpack-bottle.vercel.app/destinazioni/{slug}",
                    font_size=7.5, color=MUTED, font_name="Courier")

    add_textbox(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
                "Each destination has its own page: 3-day itinerary, neighbourhoods, food, FAQ, weather, budget tip — Italian + English.",
                font_size=10, color=INK, align=PP_ALIGN.CENTER, italic=True)

build_cities_grid_slide()

# ============= SECTION C: COOKIE CONSENT TRACKING ===============
def build_cookie_tracking_slide():
    s = section_slide("Compliance · GDPR", "Cookie consent — measured, not just enforced",
                      "Every Accept/Reject click is a tracked event in GA4")
    add_textbox(s, Inches(0.5), Inches(2.4), Inches(7), Inches(0.3),
                "WHAT FIRES", font_size=10, bold=True, color=BURNT)
    add_textbox(s, Inches(0.5), Inches(2.7), Inches(7), Inches(2.5),
                "When the user clicks Accept or Reject on the cookie banner, the site fires "
                "TWO dataLayer events:\n\n"
                "1. consent_update — sets the four Consent Mode v2 signals\n"
                "    (ad_storage, analytics_storage, ad_user_data, ad_personalization)\n\n"
                "2. consent_choice — discrete event with a single param:\n"
                "    consent_choice = 'accept' | 'reject'  ·  consent_locale = 'it' | 'en'\n\n"
                "GTM forwards consent_choice to GA4 as a key event so we can count Accept-rate "
                "vs Reject-rate as a first-class metric in Looker Studio.",
                font_size=10, color=INK)

    # Right side — code preview
    add_rect(s, Inches(8), Inches(2.4), Inches(4.8), Inches(4.0), FOREST_DARK)
    add_textbox(s, Inches(8.2), Inches(2.5), Inches(4.4), Inches(0.3),
                "components/CookieBanner.tsx", font_size=8.5, color=BURNT, font_name="Courier")
    code_lines = [
        "function decide(granted: boolean) {",
        "  localStorage.setItem('bb_consent',",
        "    granted ? 'granted' : 'denied');",
        "  pushConsent(granted);",
        "  // measured separately:",
        "  dataLayer.push({",
        "    event: 'consent_choice',",
        "    consent_choice: granted",
        "      ? 'accept' : 'reject',",
        "    consent_locale: locale,",
        "  });",
        "  setShow(false);",
        "}",
    ]
    for i, line in enumerate(code_lines):
        add_textbox(s, Inches(8.2), Inches(2.85 + i * 0.22), Inches(4.4), Inches(0.22),
                    line, font_size=8.5, color=CREAM, font_name="Courier")

    add_textbox(s, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
                "Why two events: consent_update sets state (Google's spec). consent_choice is human-readable for "
                "GA4 reporting. Both fire on the same click — instant Accept/Reject metrics.",
                font_size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

build_cookie_tracking_slide()

# ============= SECTION D: GIT REPO SCREENSHOT ===================
def build_github_slide():
    s = section_slide("Repository", "Code is public, every commit is auditable",
                      "github.com/saikanagat1117-sys/backpack-bottle")
    img = os.path.join(SHOTS, "github-repo.png")
    if os.path.exists(img):
        pi = PILImage.open(img)
        ar = pi.size[1] / pi.size[0]
        final_w = 11.5
        final_h = final_w * ar
        if final_h > 4.5:
            final_h = 4.5
            final_w = final_h / ar
        left = (13.33 - final_w) / 2
        s.shapes.add_picture(img, Inches(left), Inches(2.3),
                             width=Inches(final_w), height=Inches(final_h))
    add_textbox(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                "Configuration as code: GTM container, lead pipeline, brand tokens, briefings — "
                "every artefact lives in this repo and ships with the site.",
                font_size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

build_github_slide()

# ============= SECTION E: MEASUREMENT PAGE LIVE =================
def build_measurement_live_slide():
    s = section_slide("Live · Measurement", "Every event, documented on the live site",
                      "/measurement · KPI tree, dataLayer schema, custom dimensions, Consent Mode v2 audit")
    img = os.path.join(SHOTS, "measurement.png")
    if os.path.exists(img):
        pi = PILImage.open(img)
        ar = pi.size[1] / pi.size[0]
        final_w = 11.5
        final_h = final_w * ar
        if final_h > 4.5:
            final_h = 4.5
            final_w = final_h / ar
        left = (13.33 - final_w) / 2
        s.shapes.add_picture(img, Inches(left), Inches(2.3),
                             width=Inches(final_w), height=Inches(final_h))

build_measurement_live_slide()

# ============= SECTION F: LIVE STACK PAGE =======================
def build_stack_live_slide():
    s = section_slide("Live · Architecture", "/stack — the full stack as a live page",
                      "Same content as the visual above, browseable + shareable")
    img = os.path.join(SHOTS, "stack.png")
    if os.path.exists(img):
        pi = PILImage.open(img)
        ar = pi.size[1] / pi.size[0]
        final_w = 11.5
        final_h = final_w * ar
        if final_h > 4.5:
            final_h = 4.5
            final_w = final_h / ar
        left = (13.33 - final_w) / 2
        s.shapes.add_picture(img, Inches(left), Inches(2.3),
                             width=Inches(final_w), height=Inches(final_h))

build_stack_live_slide()

# ============= SECTION G: HOMEPAGE FULL =========================
def build_home_full_slide():
    s = section_slide("Live · Homepage", "All 11 sections of the landing flow",
                      "backpack-bottle.vercel.app · ~1.0s LCP · responsive")
    img = os.path.join(SHOTS, "home.png")
    if os.path.exists(img):
        pi = PILImage.open(img)
        ar = pi.size[1] / pi.size[0]
        # Scale tall homepage to fit
        final_h = 4.6
        final_w = final_h / ar
        if final_w > 11:
            final_w = 11
            final_h = final_w * ar
        left = (13.33 - final_w) / 2
        s.shapes.add_picture(img, Inches(left), Inches(2.3),
                             width=Inches(final_w), height=Inches(final_h))

build_home_full_slide()

# ============= SECTION H: MOBILE RESPONSIVE =====================
def build_mobile_slide():
    s = section_slide("Live · Mobile", "375px viewport — responsive proof",
                      "Mobile-first by default · 68% of Italian millennial sessions on mobile")
    img = os.path.join(SHOTS, "home-mobile.png")
    if os.path.exists(img):
        pi = PILImage.open(img)
        ar = pi.size[1] / pi.size[0]
        final_h = 4.6
        final_w = final_h / ar
        left = (13.33 - final_w) / 2
        s.shapes.add_picture(img, Inches(left), Inches(2.3),
                             width=Inches(final_w), height=Inches(final_h))
    add_textbox(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
                "Tested at 375 × 812 (iPhone 13/14/15 viewport). Same dataLayer events fire — mobile + desktop = unified analytics.",
                font_size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

build_mobile_slide()

# ============= SECTION I: GRAZIE SCREENSHOT =====================
def build_grazie_slide():
    s = section_slide("Live · Conversion", "/grazie — server-confirmed conversion",
                      "Thank-you page · BB50 coupon revealed · noindex · fires three Key Events")
    img = os.path.join(SHOTS, "grazie.png")
    if os.path.exists(img):
        pi = PILImage.open(img)
        ar = pi.size[1] / pi.size[0]
        final_w = 11
        final_h = final_w * ar
        if final_h > 4.4:
            final_h = 4.4
            final_w = final_h / ar
        left = (13.33 - final_w) / 2
        s.shapes.add_picture(img, Inches(left), Inches(2.3),
                             width=Inches(final_w), height=Inches(final_h))

build_grazie_slide()

# ============= SECTION J: PRIVACY SCREENSHOT ====================
def build_privacy_slide():
    s = section_slide("Live · Compliance", "/privacy — GDPR + cookie + coupon T&Cs",
                      "Seven sections · Italian + English · audit-ready")
    img = os.path.join(SHOTS, "privacy.png")
    if os.path.exists(img):
        pi = PILImage.open(img)
        ar = pi.size[1] / pi.size[0]
        final_w = 11
        final_h = final_w * ar
        if final_h > 4.4:
            final_h = 4.4
            final_w = final_h / ar
        left = (13.33 - final_w) / 2
        s.shapes.add_picture(img, Inches(left), Inches(2.3),
                             width=Inches(final_w), height=Inches(final_h))

build_privacy_slide()

# ============= SECTION K: BLOG ==================================
def build_blog_slide():
    s = section_slide("Live · Content", "Blog — 3 SEO-optimised posts",
                      "Italian + English · long-form · embedded internal links to destinations")
    img = os.path.join(SHOTS, "blog.png")
    if os.path.exists(img):
        pi = PILImage.open(img)
        ar = pi.size[1] / pi.size[0]
        final_w = 11
        final_h = final_w * ar
        if final_h > 4.4:
            final_h = 4.4
            final_w = final_h / ar
        left = (13.33 - final_w) / 2
        s.shapes.add_picture(img, Inches(left), Inches(2.3),
                             width=Inches(final_w), height=Inches(final_h))

build_blog_slide()

# ============= SECTION L: TEAM & ROLES ==========================
def build_team_slide():
    s = section_slide("Group 3 · Team", "Six roles, one shipped campaign",
                      "Each role owns a slice of the stack")
    roles = [
        ("01", "Sai Prathyaksh Kanagat", "CRO Specialist + Lead", "Landing page, brand, conversion."),
        ("02", "Marica Motta", "MarTech Specialist", "GTM, GA4 events, dataLayer."),
        ("03", "Alberto Faggiotto", "Compliance Specialist", "Cookie banner, Consent Mode v2, GDPR."),
        ("04", "Cecilia Gullett", "Email Marketing Specialist", "Brevo, 8-email welcome series."),
        ("05", "Marilu Stevens", "Marketing Automation", "Form → Brevo bridge, Meta CAPI, UTM persistence."),
        ("06", "Amelia Ruben", "Campaign Analyst", "KPI tree, custom dims, Looker Studio."),
    ]
    y = 2.4
    for i, (num, name, role, scope) in enumerate(roles):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.15
        yy = y + row * 1.5
        add_rect(s, Inches(x), Inches(yy), Inches(5.85), Inches(1.35), CREAM_DARK)
        add_textbox(s, Inches(x + 0.2), Inches(yy + 0.1), Inches(0.6), Inches(0.3),
                    num, font_size=11, bold=True, color=BURNT)
        add_textbox(s, Inches(x + 0.7), Inches(yy + 0.08), Inches(5.0), Inches(0.4),
                    name, font_size=14, bold=True, color=FOREST, font_name="Times-Bold")
        add_textbox(s, Inches(x + 0.7), Inches(yy + 0.45), Inches(5.0), Inches(0.3),
                    role, font_size=10, bold=True, color=BURNT_DARK)
        add_textbox(s, Inches(x + 0.7), Inches(yy + 0.75), Inches(5.0), Inches(0.5),
                    scope, font_size=10, color=INK)

build_team_slide()

# ============= SECTION M: A/B TEST =============================
def build_ab_slide():
    s = section_slide("Optimization · 01", "A/B test plan",
                      "What we'd run on day 1 to start learning fast")
    add_textbox(s, Inches(0.5), Inches(2.4), Inches(2.2), Inches(0.3),
                "HYPOTHESIS", font_size=10, bold=True, color=BURNT)
    add_textbox(s, Inches(0.5), Inches(2.7), Inches(12), Inches(1.0),
                "An emotional hero CTA ('Discover your weekend') will out-convert a transactional one "
                "('Get your €50 coupon') for cold audiences in awareness phase. Italian millennials respond to "
                "aspiration over discount when they're not yet brand-aware.",
                font_size=11, color=INK)

    # Variant A
    add_rect(s, Inches(0.5), Inches(4.0), Inches(6.0), Inches(1.4), CREAM_DARK)
    add_textbox(s, Inches(0.7), Inches(4.1), Inches(5.5), Inches(0.3),
                "VARIANT A · CONTROL", font_size=9, bold=True, color=BURNT)
    add_textbox(s, Inches(0.7), Inches(4.4), Inches(5.5), Inches(0.5),
                '"Get your €50 coupon →"', font_size=15, bold=True, color=FOREST, font_name="Times-Bold")
    add_textbox(s, Inches(0.7), Inches(4.95), Inches(5.5), Inches(0.4),
                "Transactional · direct · current production copy", font_size=9, color=MUTED)

    # Variant B
    add_rect(s, Inches(6.83), Inches(4.0), Inches(6.0), Inches(1.4), CREAM_DARK)
    add_textbox(s, Inches(7.03), Inches(4.1), Inches(5.5), Inches(0.3),
                "VARIANT B · CHALLENGER", font_size=9, bold=True, color=BURNT)
    add_textbox(s, Inches(7.03), Inches(4.4), Inches(5.5), Inches(0.5),
                '"Discover your next weekend →"', font_size=15, bold=True, color=FOREST, font_name="Times-Bold")
    add_textbox(s, Inches(7.03), Inches(4.95), Inches(5.5), Inches(0.4),
                "Emotional · aspirational · for cold awareness audiences", font_size=9, color=MUTED)

    # Decision criteria
    add_rect(s, Inches(0.5), Inches(5.6), Inches(12.33), Inches(1.4), FOREST)
    add_textbox(s, Inches(0.7), Inches(5.7), Inches(3), Inches(0.3),
                "DECISION CRITERIA", font_size=10, bold=True, color=BURNT)
    rows = [
        ("Primary metric", "form_submit conversion rate"),
        ("Sample size", "~750 sessions per variant (3% baseline · 20% MDE · 95% confidence)"),
        ("Test duration", "10–14 days minimum (covers weekday + weekend behaviour)"),
        ("Decision rule", "Ship winner if uplift ≥15% with p<0.05; else keep control"),
    ]
    for i, (k, v) in enumerate(rows):
        add_textbox(s, Inches(0.7), Inches(6.0 + i*0.23), Inches(2.5), Inches(0.22),
                    k, font_size=9, bold=True, color=CREAM)
        add_textbox(s, Inches(3.3), Inches(6.0 + i*0.23), Inches(9.3), Inches(0.22),
                    v, font_size=9, color=CREAM)

build_ab_slide()

# ============= SECTION N: COMPETITOR BENCHMARK ==================
def build_competitor_slide():
    s = section_slide("Optimization · 02", "Competitor benchmark",
                      "Italian travel sector · CPL, CTR, positioning · what we're up against")
    hdr = ["Competitor", "Positioning", "CPL", "CTR", "Lead magnet", "Our edge"]
    col_widths = [Inches(2.0), Inches(2.5), Inches(1.0), Inches(1.0), Inches(2.5), Inches(2.5)]
    xs = [Inches(0.5)]
    for w in col_widths[:-1]:
        xs.append(xs[-1] + w)
    y0 = Inches(2.5)
    add_rect(s, Inches(0.5), y0, Inches(12), Inches(0.4), FOREST)
    for i, h in enumerate(hdr):
        add_textbox(s, xs[i] + Inches(0.1), y0 + Inches(0.1), col_widths[i], Inches(0.3),
                    h.upper(), font_size=9, bold=True, color=CREAM)

    rows = [
        ("Volagratis", "Cheapest flights, no curation", "€11.20", "1.8%", "Newsletter signup", "We curate; they don't"),
        ("eDreams Prime", "Membership flight discounts", "€9.50", "2.0%", "Free Prime trial", "No subscription friction"),
        ("Lastminute.com", "Last-minute travel deals", "€8.30", "2.4%", "App download", "We're early-booking, not last-minute"),
        ("BACKPACK & BOTTLE", "Curated · transparent · early-booking", "€6.67", "2.7% (proj)", "BB50 €50 coupon", "First IT to combine all 3"),
    ]
    for r, row in enumerate(rows):
        yy = y0 + Inches(0.4 + r*0.55)
        bg = CREAM if r < 3 else BURNT
        add_rect(s, Inches(0.5), yy, Inches(12), Inches(0.55), bg)
        for i, val in enumerate(row):
            text_color = CREAM if r == 3 else INK
            bold = (i == 0) or (r == 3)
            add_textbox(s, xs[i] + Inches(0.1), yy + Inches(0.16),
                        col_widths[i], Inches(0.3),
                        val, font_size=10, bold=bold, color=text_color)
    add_textbox(s, Inches(0.5), Inches(5.7), Inches(12), Inches(1.3),
                "Source: WordStream travel benchmarks 2025, SimilarWeb traffic data, Italy-localised CPC estimates from "
                "Google Keyword Planner. Our projected CPL of €6.67 is ~30% better than the closest competitor (Lastminute) — "
                "driven by tighter audience targeting (Italian millennials only) and a higher-value lead magnet (€50 coupon).",
                font_size=10, color=MUTED)

build_competitor_slide()

# ============= SECTION O: ACCESSIBILITY =========================
def build_a11y_slide():
    s = section_slide("Optimization · 03", "Accessibility audit",
                      "WCAG 2.1 AA review · 6 PASS / 3 PARTIAL · documented fixes")
    a11y = [
        ("✓", "Color contrast", "Forest #1F3A2E on cream = 12.4:1 (AAA). Burnt on cream = 4.6:1 (AA).", "PASS"),
        ("✓", "Heading hierarchy", "Single <h1> per page; h2/h3 cascade. Nav uses semantic <nav>.", "PASS"),
        ("✓", "Alt text", "All images have descriptive alt. Decorative ones use aria-hidden=true.", "PASS"),
        ("◐", "Form labels", "All 5 fields labelled. Could improve: explicit aria-required=true.", "PARTIAL"),
        ("◐", "Keyboard nav", "Tab traversal works. Sticky CTA may need skip-link. Modal traps focus correctly.", "PARTIAL"),
        ("◐", "Cookie banner", "Dismissible via Escape. Should announce as aria-live=polite. Fix in v2.", "PARTIAL"),
        ("✓", "Focus indicators", "Default focus rings preserved. Visible on all interactive elements.", "PASS"),
        ("✓", "Lang attribute", "<html lang='it'> set; alternate hreflang for /?lang=en in metadata.", "PASS"),
    ]
    y0 = 2.5
    for i, (icon, item, desc, status) in enumerate(a11y):
        yy = y0 + i * 0.5
        color = FOREST if status == "PASS" else BURNT
        add_textbox(s, Inches(0.5), Inches(yy), Inches(0.5), Inches(0.4),
                    icon, font_size=14, bold=True, color=color)
        add_textbox(s, Inches(1.0), Inches(yy + 0.05), Inches(2.5), Inches(0.4),
                    item, font_size=11, bold=True, color=FOREST)
        add_textbox(s, Inches(3.5), Inches(yy + 0.07), Inches(8), Inches(0.4),
                    desc, font_size=9, color=INK)
        add_textbox(s, Inches(11.7), Inches(yy + 0.08), Inches(1.0), Inches(0.4),
                    status, font_size=9, bold=True, color=color, align=PP_ALIGN.RIGHT)

    add_textbox(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
                "Tooling: Manual code review + axe-core DevTools spot-check. Documented fixes shipped to v2 backlog.",
                font_size=9, color=MUTED, italic=True)

build_a11y_slide()

# =================================================================
# SAVE
# =================================================================
prs.save(OUT)
print(f"\nSaved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
