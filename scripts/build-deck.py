#!/usr/bin/env python3
"""Backpack & Bottle — Group 3 — Digital Platforms Lab campaign deck.
Generates Backpack-Bottle-Campaign.pptx covering brand, strategy, tracking, KPIs, media plan."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# Brand
FOREST = RGBColor(0x1F, 0x3A, 0x2E)
FOREST_DK = RGBColor(0x14, 0x26, 0x20)
CREAM = RGBColor(0xF5, 0xEF, 0xE6)
CREAM_DK = RGBColor(0xE8, 0xDF, 0xD1)
BURNT = RGBColor(0xD9, 0x76, 0x42)
BURNT_DK = RGBColor(0xB8, 0x5E, 0x30)
INK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ---- helpers ----

def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK, font="Inter", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill, line=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    if line:
        r.line.color.rgb = line
        r.line.width = Pt(1)
    else:
        r.line.fill.background()
    r.shadow.inherit = False
    return r


def add_rounded(slide, x, y, w, h, fill, line=None):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    if line:
        r.line.color.rgb = line
    else:
        r.line.fill.background()
    r.shadow.inherit = False
    return r


def header(slide, title, subtitle=None, dark=False):
    """Standard slide header"""
    bg_color = CREAM
    text_color = FOREST
    accent_color = BURNT
    if dark:
        bg_color = FOREST
        text_color = CREAM
        accent_color = BURNT
    add_bg(slide, bg_color)
    # eyebrow bar
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(0.4), Pt(2), accent_color)
    add_text(slide, Inches(1.1), Inches(0.45), Inches(11), Inches(0.4),
             "GROUP 3 · DIGITAL PLATFORMS LAB",
             size=9, bold=True, color=accent_color, font="Inter")
    add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(1.0),
             title, size=32, bold=True, color=text_color, font="Georgia")
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
                 subtitle, size=14, color=text_color, font="Inter")
    return text_color, accent_color


def footer(slide, page_num, total, dark=False):
    color = CREAM if dark else FOREST
    add_text(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
             "backpack & bottle  ·  Group 3", size=9, color=color, font="Inter")
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
             f"{page_num} / {total}", size=9, color=color, font="Inter", align=PP_ALIGN.RIGHT)


def add_table(slide, x, y, w, h, headers, rows, *, header_fill=FOREST, header_text=CREAM, body_text=INK, border=CREAM_DK, font_size=10):
    nrows = len(rows) + 1
    ncols = len(headers)
    tbl_shape = slide.shapes.add_table(nrows, ncols, x, y, w, h)
    tbl = tbl_shape.table
    # header
    for j, htext in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        tf.text = ""
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = htext
        r.font.bold = True
        r.font.size = Pt(font_size)
        r.font.name = "Inter"
        r.font.color.rgb = header_text
    # body
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else CREAM
            tf = cell.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            tf.word_wrap = True
            tf.text = ""
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(font_size)
            r.font.name = "Inter"
            r.font.color.rgb = body_text
    return tbl


# ---- slide builders ----

slides_data = []  # populated via decorator

def slide(fn):
    slides_data.append(fn)
    return fn


@slide
def s_cover(s, idx, total):
    add_bg(s, FOREST)
    # Decorative
    add_rect(s, Inches(0.6), Inches(0.55), Inches(0.4), Pt(2), BURNT)
    add_text(s, Inches(1.1), Inches(0.45), Inches(11), Inches(0.4),
             "GROUP 3 · DIGITAL PLATFORMS LAB · BOLOGNA BUSINESS SCHOOL",
             size=10, bold=True, color=BURNT, font="Inter")
    add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(2.5),
             "backpack & bottle.", size=72, bold=True, color=CREAM, font="Georgia")
    add_text(s, Inches(0.6), Inches(3.6), Inches(11), Inches(1.4),
             "Curated European city breaks. Transparent prices.\n€50 off your first booking.",
             size=22, color=CREAM, font="Georgia")
    add_text(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.4),
             "Campaign brief · brand identity · landing page · tracking · media plan",
             size=12, color=BURNT, font="Inter")
    add_text(s, Inches(0.6), Inches(6.8), Inches(8), Inches(0.5),
             "Sai Prathyaksh Kanagat  ·  Group 3", size=12, color=CREAM, font="Inter")
    add_text(s, Inches(8.5), Inches(6.8), Inches(4.4), Inches(0.5),
             "Campaign window: 5 May – 15 June 2026",
             size=12, color=CREAM, font="Inter", align=PP_ALIGN.RIGHT)


@slide
def s_agenda(s, idx, total):
    header(s, "Agenda", "What we'll walk through")
    items = [
        ("01", "The brief", "What we were asked to do"),
        ("02", "Strategic decisions", "Who, what, when, how much"),
        ("03", "The audience", "Italian millennials, persona + pain points"),
        ("04", "Brand identity", "Positioning, voice, visual system"),
        ("05", "Offer & coupon", "5 destinations, BB50 mechanic"),
        ("06", "Landing page", "Architecture + walkthrough"),
        ("07", "Measurement plan", "GTM, GA4, Meta Pixel, Consent Mode v2"),
        ("08", "Media plan & budget", "60/40 Meta/Google, €2,000"),
        ("09", "Timeline", "Pre-launch → 6-week run → analysis"),
        ("10", "Risks & next steps", "Assumptions and what we'd do live"),
    ]
    y = Inches(2.4)
    for i, (n, t, sub) in enumerate(items):
        col = i // 5
        row = i % 5
        x = Inches(0.6 + col * 6.3)
        yy = y + Inches(row * 0.85)
        add_text(s, x, yy, Inches(0.6), Inches(0.4), n, size=18, bold=True, color=BURNT, font="Georgia")
        add_text(s, x + Inches(0.7), yy, Inches(5.2), Inches(0.4), t, size=14, bold=True, color=FOREST, font="Inter")
        add_text(s, x + Inches(0.7), yy + Inches(0.35), Inches(5.2), Inches(0.4), sub, size=11, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_brief(s, idx, total):
    header(s, "The brief", "Group 3 — Backpack & Bottle")
    cols = [
        ("Business objective", "Promote travel deals under the \"Backpack & Bottle\" brand and generate qualified leads interested in affordable, experience-driven trips."),
        ("Marketing objective", "Make the offer feel attractive, accessible, and time-sensitive — combining inspiration with a clear incentive to request a coupon and book early."),
        ("What was open", "Audience, offer architecture, value prop sharpening, media mix, budget, timing, KPI targets — all left to the team to define and justify."),
    ]
    y = Inches(2.6)
    for i, (h, body) in enumerate(cols):
        x = Inches(0.6 + i * 4.2)
        add_rounded(s, x, y, Inches(4.0), Inches(3.6), WHITE)
        add_rect(s, x, y, Inches(4.0), Pt(3), BURNT)
        add_text(s, x + Inches(0.3), y + Inches(0.3), Inches(3.4), Inches(0.4),
                 h.upper(), size=10, bold=True, color=BURNT, font="Inter")
        add_text(s, x + Inches(0.3), y + Inches(0.8), Inches(3.4), Inches(2.7),
                 body, size=12, color=FOREST, font="Inter")
    add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
             "Mandatory KPIs: LP visits · coupon downloads · lead form submissions · CVR · CPL · email list growth.",
             size=11, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_decisions(s, idx, total):
    header(s, "Strategic decisions at a glance")
    items = [
        ("Audience", "Italian millennials 25–39, urban, €25–45k income"),
        ("Offer", "Curated European city breaks · 5 destinations"),
        ("Coupon", "€50 off · code BB50 · valid 90 days · min €300"),
        ("Positioning", "\"Skip the research fatigue\""),
        ("Strategic focus", "Lead acquisition for email nurture (long booking cycle)"),
        ("Budget", "€2,000 paid media · 60% Meta · 40% Google"),
        ("Window", "5 May – 15 June 2026 (6 weeks active)"),
        ("Targets", "5,000 visits → 400 coupons → 300 leads · €6.67 CPL"),
        ("Stack", "Next.js + Vercel · GTM · GA4 · Meta Pixel · Looker Studio"),
    ]
    y = Inches(2.4)
    for i, (h, body) in enumerate(items):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.15)
        yy = y + Inches(row * 1.4)
        add_rounded(s, x, yy, Inches(3.95), Inches(1.25), WHITE)
        add_text(s, x + Inches(0.25), yy + Inches(0.18), Inches(3.5), Inches(0.35),
                 h.upper(), size=9, bold=True, color=BURNT, font="Inter")
        add_text(s, x + Inches(0.25), yy + Inches(0.55), Inches(3.5), Inches(0.7),
                 body, size=12, bold=True, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_audience_who(s, idx, total):
    header(s, "Who we're talking to", "Italian millennials 25–39 · urban · digitally native")
    add_rounded(s, Inches(0.6), Inches(2.5), Inches(5.8), Inches(4.2), WHITE)
    add_text(s, Inches(0.9), Inches(2.7), Inches(5.4), Inches(0.4),
             "DEMOGRAPHICS", size=10, bold=True, color=BURNT, font="Inter")
    demo = [
        "Age: 25–39",
        "Location: Milan · Rome · Bologna · Florence · Turin",
        "Income: €25–45k/year",
        "Education: University degree",
        "Employment: Digital, creative, service sectors",
        "Relationship: Single or couples without children",
        "Language: Italian (primary), English (secondary)",
    ]
    add_text(s, Inches(0.9), Inches(3.1), Inches(5.4), Inches(3.5),
             "\n".join(demo), size=12, color=FOREST, font="Inter")
    # Right column
    add_rounded(s, Inches(6.7), Inches(2.5), Inches(6.0), Inches(4.2), FOREST)
    add_text(s, Inches(7.0), Inches(2.7), Inches(5.6), Inches(0.4),
             "TRAVEL MOTIVATIONS & BEHAVIOUR", size=10, bold=True, color=BURNT, font="Inter")
    motiv = [
        "• Authentic local life — not package tourism",
        "• Quality time with friends/family (38% motivation)",
        "• Food & wine experiences (36%)",
        "• Shareable Instagram moments",
        "• Books 2–4 months ahead, takes 3–4 trips/year",
        "• 70% inspired by destinations seen on social",
        "• Mobile for inspiration · desktop for booking",
        "• Avoids August (Ferragosto), prefers Apr–Jun & Sep–Oct",
    ]
    add_text(s, Inches(7.0), Inches(3.1), Inches(5.6), Inches(3.5),
             "\n".join(motiv), size=12, color=CREAM, font="Inter")
    footer(s, idx, total)


@slide
def s_audience_pains(s, idx, total):
    header(s, "What blocks them from booking", "Five pain points the campaign solves")
    pains = [
        ("Research overwhelm", "Booking.com shows 500 hotels — paralysing. We show 3–5 curated options per city."),
        ("Price transparency", "Frustrated by prices that jump at checkout. We show all-in pricing upfront."),
        ("Trust deficit", "Burned by Volagratis / eDreams refund delays. We position transparency as the brand."),
        ("FOMO on the \"best\" deal", "Constantly worry they're missing a better price. €50 off is a real, named discount — not a fake countdown."),
        ("Time scarcity", "Want to travel more but exhausted by planning. Pre-packaged itineraries cut planning time to 5 minutes."),
    ]
    y = Inches(2.4)
    for i, (h, body) in enumerate(pains):
        yy = y + Inches(i * 0.9)
        add_text(s, Inches(0.6), yy, Inches(0.6), Inches(0.6),
                 f"0{i+1}", size=22, bold=True, color=BURNT, font="Georgia")
        add_text(s, Inches(1.4), yy, Inches(3.2), Inches(0.4),
                 h, size=14, bold=True, color=FOREST, font="Inter")
        add_text(s, Inches(4.7), yy, Inches(8.0), Inches(0.7),
                 body, size=12, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_brand_position(s, idx, total):
    header(s, "Positioning", "How we differ from every competitor in the audience's head", dark=True)
    add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.6),
             "\"Skip the research fatigue.\"", size=36, bold=True, color=CREAM, font="Georgia")
    add_text(s, Inches(0.6), Inches(3.2), Inches(12), Inches(0.5),
             "European city breaks at insider prices — early access deals that fit your budget and deliver experiences worth sharing.",
             size=14, color=CREAM, font="Inter")
    rows = [
        ("Booking.com / Expedia", "Compete on breadth — thousands of options", "We compete on curation — best 3–5 per city"),
        ("Volagratis / Lastminute", "Compete on price (but with trust issues)", "We compete on transparent value + trust"),
        ("Luxury tour operators", "High-touch concierge service", "Affordable authenticity"),
        ("DIY booking", "Maximum control, maximum effort", "5-minute booking, time saved"),
    ]
    add_table(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.3),
              ["Competitor", "How they compete", "How we win"], rows,
              header_fill=BURNT, header_text=CREAM, body_text=FOREST,
              font_size=11)
    footer(s, idx, total, dark=True)


@slide
def s_brand_visual(s, idx, total):
    header(s, "Brand identity", "Modern Wanderer — warm, editorial, trustworthy")
    # Palette
    palette = [
        ("Forest", "#1F3A2E", FOREST),
        ("Forest dark", "#142620", FOREST_DK),
        ("Cream", "#F5EFE6", CREAM),
        ("Cream dark", "#E8DFD1", CREAM_DK),
        ("Burnt orange", "#D97642", BURNT),
        ("Burnt dark", "#B85E30", BURNT_DK),
    ]
    add_text(s, Inches(0.6), Inches(2.4), Inches(6), Inches(0.4),
             "PALETTE", size=10, bold=True, color=BURNT, font="Inter")
    for i, (n, hx, c) in enumerate(palette):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 2.0)
        y = Inches(2.8 + row * 1.5)
        add_rect(s, x, y, Inches(1.8), Inches(1.0), c)
        add_text(s, x, y + Inches(1.05), Inches(1.8), Inches(0.3), n, size=10, bold=True, color=FOREST, font="Inter")
        add_text(s, x, y + Inches(1.3), Inches(1.8), Inches(0.3), hx, size=9, color=FOREST, font="Inter")
    # Type
    add_text(s, Inches(7.0), Inches(2.4), Inches(6), Inches(0.4),
             "TYPOGRAPHY", size=10, bold=True, color=BURNT, font="Inter")
    add_text(s, Inches(7.0), Inches(2.8), Inches(6), Inches(0.6),
             "Fraunces", size=36, bold=True, color=FOREST, font="Georgia")
    add_text(s, Inches(7.0), Inches(3.5), Inches(6), Inches(0.4),
             "Display serif — used for headlines, wordmark, prices",
             size=11, color=FOREST, font="Inter")
    add_text(s, Inches(7.0), Inches(4.2), Inches(6), Inches(0.6),
             "Inter", size=32, bold=True, color=FOREST, font="Inter")
    add_text(s, Inches(7.0), Inches(4.9), Inches(6), Inches(0.4),
             "Sans — used for body, UI, captions",
             size=11, color=FOREST, font="Inter")
    # Voice
    add_text(s, Inches(7.0), Inches(5.6), Inches(6), Inches(0.4),
             "VOICE", size=10, bold=True, color=BURNT, font="Inter")
    add_text(s, Inches(7.0), Inches(5.95), Inches(6), Inches(1.2),
             "Direct. Warm. Anti-tourist-trap. Editorial without being pretentious. We say what we don't do as much as what we do.",
             size=12, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_value_pillars(s, idx, total):
    header(s, "The promise + 4 selling points")
    add_rounded(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(1.2), FOREST)
    add_text(s, Inches(0.9), Inches(2.65), Inches(11.5), Inches(0.4),
             "PROMISE", size=10, bold=True, color=BURNT, font="Inter")
    add_text(s, Inches(0.9), Inches(2.95), Inches(11.5), Inches(0.7),
             "We do the hard work of curation so you don't have to spend hours comparing options. We pick the best 3–5 packages per destination, negotiate early-booking rates, and present them transparently. You book in 5 minutes and go create memories.",
             size=12, color=CREAM, font="Inter")
    pillars = [
        ("Curated, not overwhelming", "3–5 great options per city, not 500. Every hotel pre-vetted for location, reviews, value."),
        ("Insider early-booking pricing", "€50 off through partnership rates. Real discount, transparent all-in pricing."),
        ("Authentic local experiences", "Neighbourhoods locals actually live in. Real cafés and markets, not tourist traps."),
        ("Worry-free planning", "Flights + hotel + itinerary coordinated. Flexible dates. Email support if questions arise."),
    ]
    y = Inches(4.0)
    for i, (h, body) in enumerate(pillars):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.1)
        yy = y + Inches(row * 1.4)
        add_rounded(s, x, yy, Inches(5.95), Inches(1.25), WHITE)
        add_rect(s, x, yy, Inches(0.1), Inches(1.25), BURNT)
        add_text(s, x + Inches(0.3), yy + Inches(0.18), Inches(5.5), Inches(0.4),
                 h, size=14, bold=True, color=FOREST, font="Inter")
        add_text(s, x + Inches(0.3), yy + Inches(0.6), Inches(5.5), Inches(0.7),
                 body, size=11, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_offer(s, idx, total):
    header(s, "The offer", "5 curated European city breaks · 1.5–3 hour flights from Italy")
    dests = [
        ("Rome", "Italy", "€280 · 2 nights", "Classic"),
        ("Barcelona", "Spain", "€320 · 3 nights", "Coast & city"),
        ("Amsterdam", "Netherlands", "€365 · 3 nights", "Canals & design"),
        ("Lisbon", "Portugal", "€345 · 3 nights", "Atlantic"),
        ("Prague", "Czech Republic", "€295 · 2 nights", "Historic"),
    ]
    y = Inches(2.5)
    for i, (city, country, price, tag) in enumerate(dests):
        x = Inches(0.6 + i * 2.55)
        add_rounded(s, x, y, Inches(2.4), Inches(3.4), FOREST)
        add_rect(s, x + Inches(0.2), y + Inches(0.2), Inches(2.0), Inches(0.3), BURNT)
        add_text(s, x + Inches(0.25), y + Inches(0.22), Inches(1.95), Inches(0.3),
                 tag.upper(), size=9, bold=True, color=CREAM, font="Inter")
        add_text(s, x + Inches(0.2), y + Inches(0.7), Inches(2.0), Inches(0.4),
                 country.upper(), size=9, color=BURNT, font="Inter")
        add_text(s, x + Inches(0.2), y + Inches(1.05), Inches(2.0), Inches(0.7),
                 city, size=24, bold=True, color=CREAM, font="Georgia")
        add_text(s, x + Inches(0.2), y + Inches(2.0), Inches(2.0), Inches(0.4),
                 "FROM", size=9, color=BURNT, font="Inter")
        add_text(s, x + Inches(0.2), y + Inches(2.3), Inches(2.0), Inches(0.6),
                 price, size=14, bold=True, color=CREAM, font="Inter")
    add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.5),
             "Each package includes: round-trip flight from Italian airports, 3–4 star boutique hotel, local guide & itinerary suggestions. Optional add-ons: food tours, bike rentals, museum passes.",
             size=11, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_coupon(s, idx, total):
    header(s, "Coupon mechanic", "BB50 — built to drive lead capture, not race-to-the-bottom pricing", dark=True)
    add_rounded(s, Inches(0.6), Inches(2.5), Inches(6.0), Inches(4.2), CREAM)
    add_text(s, Inches(0.9), Inches(2.7), Inches(5.4), Inches(0.4),
             "MECHANIC", size=10, bold=True, color=BURNT, font="Inter")
    add_text(s, Inches(0.9), Inches(3.1), Inches(5.4), Inches(1.5),
             "€50", size=72, bold=True, color=BURNT, font="Georgia")
    add_text(s, Inches(0.9), Inches(4.7), Inches(5.4), Inches(0.5),
             "off bookings from €300", size=14, color=FOREST, font="Inter")
    add_text(s, Inches(0.9), Inches(5.2), Inches(5.4), Inches(0.4),
             "Code: BB50  ·  Valid 90 days  ·  One per customer",
             size=12, bold=True, color=FOREST, font="Inter")
    add_text(s, Inches(0.9), Inches(5.65), Inches(5.4), Inches(0.9),
             "Travel dates 1 Jun – 31 Dec 2026.\nExcl. 23–31 Dec, 31 Dec – 2 Jan, 10–20 Aug.\nNot combinable with other offers.",
             size=10, color=FOREST, font="Inter")
    add_rounded(s, Inches(6.9), Inches(2.5), Inches(5.8), Inches(4.2), CREAM)
    add_text(s, Inches(7.2), Inches(2.7), Inches(5.2), Inches(0.4),
             "WHY THIS DESIGN", size=10, bold=True, color=BURNT, font="Inter")
    why = [
        "€300 minimum: realistic floor for flight + hotel city break.",
        "€50 = 16.7% — meaningful but doesn't devalue the brand.",
        "90-day validity: matches Italian booking cycle (2–4 months ahead) without inducing urgency paralysis.",
        "Required email + GDPR consent → coupon = the lead-magnet logic.",
        "One-per-customer prevents stockpiling, supports clean reporting.",
    ]
    add_text(s, Inches(7.2), Inches(3.1), Inches(5.4), Inches(3.5),
             "\n\n".join(why), size=11, color=FOREST, font="Inter")
    footer(s, idx, total, dark=True)


@slide
def s_funnel(s, idx, total):
    header(s, "The funnel", "From paid impression to booking — and where each KPI sits")
    stages = [
        ("Impressions", "150–180k", "Meta + Google Ads"),
        ("Clicks", "1.5–2k", "CTR ~1%"),
        ("Landing visits", "5,000", "KPI #1"),
        ("Coupon downloads", "400", "KPI #2 · 8% CVR"),
        ("Email leads", "300", "KPI #3 · 6% CVR"),
        ("Bookings", "~30", "10% nurture conv."),
    ]
    y = Inches(2.7)
    for i, (n, val, sub) in enumerate(stages):
        x = Inches(0.6 + i * 2.1)
        w = Inches(2.0)
        # Funnel block — narrows visually with each stage
        h_pct = 1.0 - (i * 0.08)
        h = Inches(2.0 * h_pct)
        block_y = y + Inches((2.0 - 2.0 * h_pct) / 2)
        color = FOREST if i < 3 else BURNT
        add_rounded(s, x, block_y, w, h, color)
        add_text(s, x, block_y + Inches(0.15), w, Inches(0.4),
                 n.upper(), size=9, bold=True, color=CREAM, font="Inter", align=PP_ALIGN.CENTER)
        add_text(s, x, block_y + Inches(0.5), w, Inches(0.7),
                 val, size=20, bold=True, color=CREAM, font="Georgia", align=PP_ALIGN.CENTER)
        add_text(s, x, block_y + Inches(1.2), w, Inches(0.4),
                 sub, size=9, color=CREAM, font="Inter", align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            add_text(s, x + Inches(1.95), y + Inches(0.95), Inches(0.2), Inches(0.4),
                     "→", size=18, bold=True, color=BURNT, font="Inter", align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
             "€2,000 spend ÷ 300 leads = €6.67 CPL · ROI tracked over 90 days as nurtured leads convert to bookings.",
             size=12, bold=True, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_lp_overview(s, idx, total):
    header(s, "Landing page architecture", "Single-purpose page optimised for one action: download the coupon")
    sections = [
        ("Sticky nav", "Wordmark · Destinations · How it works · Blog · Measurement · EN/IT toggle · Coupon CTA"),
        ("Hero", "Eyebrow tag · headline (\"Skip the research fatigue\") · subheadline · dual CTA · trust line"),
        ("Value pillars", "4 pillars on dark forest background — curation, transparency, authentic, worry-free"),
        ("Destinations grid", "5 cards · image · country · city · from-price · clickable to scroll to form"),
        ("How it works", "3 steps · claim coupon → pick city → book and go"),
        ("Social proof", "4 stats: 500+ hotels vetted · 5 min to book · 0 hidden fees · €50 off"),
        ("Lead capture", "Form: name · email · departure city · destination interest · GDPR consent · CTA"),
        ("Footer", "Wordmark · navigation · privacy · terms · academic project credit"),
    ]
    y = Inches(2.4)
    for i, (h, body) in enumerate(sections):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.1)
        yy = y + Inches(row * 1.05)
        add_rect(s, x, yy + Inches(0.05), Inches(0.05), Inches(0.85), BURNT)
        add_text(s, x + Inches(0.2), yy, Inches(5.7), Inches(0.4),
                 h, size=12, bold=True, color=FOREST, font="Inter")
        add_text(s, x + Inches(0.2), yy + Inches(0.35), Inches(5.7), Inches(0.7),
                 body, size=10, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_lp_pages(s, idx, total):
    header(s, "Site map", "Every route the campaign serves")
    routes = [
        ("/", "Italian primary landing page (EN toggle)"),
        ("/grazie", "Thank-you page · BB50 displayed · PDF download · GA4 conversion event"),
        ("/measurement", "Public measurement plan (this deck, on the live site)"),
        ("/blog", "SEO hub — 3 posts on weekend itineraries, booking timing, destination comparisons"),
        ("/blog/[slug]", "Individual post pages with CTA back to coupon form"),
        ("/privacy", "GDPR privacy policy + cookie disclosures + coupon T&Cs"),
        ("/api/lead", "Server endpoint accepting form submissions, forwarding to ESP webhook"),
        ("/backpack-bottle-coupon-BB50.pdf", "Branded one-page coupon PDF served from /public"),
    ]
    add_table(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(4.2),
              ["Route", "Purpose"], routes, font_size=11)
    add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
             "Built on Next.js 14 (App Router) + TypeScript + Tailwind, deployed on Vercel.",
             size=11, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_tracking_arch(s, idx, total):
    header(s, "Tracking architecture", "One container. Every tag. Consent-aware.", dark=True)
    add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             "All marketing tags fire through a single GTM container. Each user interaction pushes a named event into window.dataLayer. GTM listens, fires GA4 + Meta Pixel + Google Ads. Consent Mode v2 blocks marketing storage until consent is granted.",
             size=12, color=CREAM, font="Inter")
    diagram = (
        "User action (click / scroll / submit / page load)\n"
        "                         │\n"
        "                         ▼\n"
        "    window.dataLayer.push({event: 'coupon_download', value: 50, ...})\n"
        "                         │\n"
        "                         ▼\n"
        " ┌────────────────────────────────────────────────────────────┐\n"
        " │  Google Tag Manager   (GTM-XXXXXXX)                         │\n"
        " │  • Trigger: Custom Event = coupon_download                  │\n"
        " │  • Variables read: coupon_code · value · currency · destination │\n"
        " │  • Consent gate: ad_storage = granted required              │\n"
        " └─────────┬──────────────────┬──────────────────┬─────────────┘\n"
        "           ▼                  ▼                  ▼\n"
        "      [GA4 event]        [Meta Lead]      [Google Ads conv.]\n"
        "       G-XXXXXX           Pixel ID         AW-XXXXXX/abcd\n"
        "           │                  │                  │\n"
        "           ▼                  ▼                  ▼\n"
        "      Looker Studio    Retargeting aud.   Conversion import"
    )
    add_rounded(s, Inches(0.6), Inches(3.3), Inches(12.1), Inches(3.6), FOREST_DK)
    add_text(s, Inches(0.9), Inches(3.4), Inches(11.5), Inches(3.5),
             diagram, size=10, color=CREAM, font="Courier New")
    footer(s, idx, total, dark=True)


@slide
def s_datalayer_events(s, idx, total):
    header(s, "dataLayer event schema", "Every custom event the site fires")
    rows = [
        ("page_view", "Auto via GA4 Enhanced Measurement", "page_location, page_title", "Funnel entry"),
        ("scroll_75", "GTM Scroll Depth trigger 75%", "percent_scrolled", "Engagement quality"),
        ("cta_click", "Hero, nav, blog CTAs", "location", "Intent signal"),
        ("destination_card_click", "Destination card click", "destination", "Custom dim destination_interest"),
        ("form_start", "First form-field focus", "form: coupon_lead", "Form abandonment"),
        ("form_submit", "POST /api/lead success", "form, destination", "GA4 lead conversion"),
        ("coupon_download", "After submit + direct PDF download", "coupon_code, value, currency", "Primary conversion"),
        ("thank_you_view", "/grazie page load", "value, currency", "Backup conversion"),
        ("consent_update", "Cookie banner choice", "ad_storage, analytics_storage, ad_user_data, ad_personalization", "Consent Mode v2"),
        ("visit_count_set", "Every page load (localStorage)", "visit_count", "User property"),
    ]
    add_table(s, Inches(0.6), Inches(2.4), Inches(12.1), Inches(4.4),
              ["Event", "Trigger", "Parameters", "Purpose"], rows, font_size=10)
    footer(s, idx, total)


@slide
def s_gtm_tags(s, idx, total):
    header(s, "GTM tag inventory", "Every tag the container ships, what it does, when it fires")
    rows = [
        ("GA4 Configuration", "All Pages", "Pageviews + base config to property G-XXXXXXXXXX"),
        ("GA4 Event — coupon_download", "Custom Event = coupon_download", "Marked as conversion in GA4"),
        ("GA4 Event — form_submit", "Custom Event = form_submit", "Lead-stage conversion"),
        ("GA4 Event — cta_click", "Custom Event = cta_click", "Engagement tracking"),
        ("GA4 Event — scroll_75", "Custom Event = scroll_75", "Engagement quality"),
        ("Meta Pixel — Base", "All Pages", "PageView for retargeting audiences"),
        ("Meta Pixel — Lead", "Custom Event = form_submit", "Optimisation event for Meta Ads"),
        ("Meta Pixel — CompleteRegistration", "Custom Event = coupon_download", "Bottom-funnel signal"),
        ("Google Ads — Lead conversion", "Custom Event = coupon_download", "Value 50 EUR; imported into Google Ads"),
        ("Google Ads — Remarketing", "All Pages", "Audience builder for retargeting"),
    ]
    add_table(s, Inches(0.6), Inches(2.4), Inches(12.1), Inches(4.4),
              ["Tag", "Trigger", "Purpose"], rows, font_size=10)
    footer(s, idx, total)


@slide
def s_consent(s, idx, total):
    header(s, "Consent Mode v2 (GDPR)", "Default-denied. Privacy-first. Audit-ready.", dark=True)
    add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             "Before GTM loads, gtag('consent', 'default', ...) sets all four signals to denied. The cookie banner asks the user. On Accept, a consent_update event flips signals to granted; pending tags execute.",
             size=12, color=CREAM, font="Inter")
    rows = [
        ("ad_storage", "denied", "granted", "Meta Pixel · Google Ads cookies"),
        ("analytics_storage", "denied", "granted", "GA4 cookies"),
        ("ad_user_data", "denied", "granted", "User-data signals to Google"),
        ("ad_personalization", "denied", "granted", "Personalised ads"),
    ]
    add_table(s, Inches(0.6), Inches(3.4), Inches(12.1), Inches(2.6),
              ["Signal", "Default state", "After accept", "Controls"], rows,
              header_fill=BURNT, header_text=CREAM, body_text=FOREST, font_size=11)
    add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
             "Even with Reject, GA4 receives consent-modeled pings (no cookies) so totals are still observable in aggregate — Google's recommended setup.",
             size=11, color=CREAM, font="Inter")
    footer(s, idx, total, dark=True)


@slide
def s_kpis(s, idx, total):
    header(s, "KPI definitions & targets")
    rows = [
        ("Landing page visits", "5,000", "€2,000 ÷ €0.40 blended CPC", "GA4 sessions"),
        ("Coupon downloads", "400 (8% CVR)", "Travel lead-magnet benchmark 5–10%", "GA4 coupon_download"),
        ("Lead form submissions", "300 (6% CVR)", "Industry standard for gated content 5–8%", "GA4 form_submit"),
        ("Conversion rate", "6–8%", "(downloads + forms) ÷ visits", "Calculated metric"),
        ("Cost per lead (CPL)", "€6.67 avg", "Travel sector typical €10–50", "Spend ÷ leads"),
        ("Email list growth", "400+", "Net new subscribers from campaign", "ESP report"),
        ("Email open rate", "35%", "Travel industry avg 30–35%", "ESP report"),
        ("Click-to-open rate", "20%", "Travel industry avg 15–25%", "ESP report"),
    ]
    add_table(s, Inches(0.6), Inches(2.4), Inches(12.1), Inches(4.4),
              ["KPI", "Target", "Rationale", "Source"], rows, font_size=11)
    footer(s, idx, total)


@slide
def s_media_overview(s, idx, total):
    header(s, "Media plan overview", "Where the €2,000 goes and why", dark=True)
    add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             "Two channels. One job each. Meta drives discovery + retargeting; Google captures active intent.",
             size=14, color=CREAM, font="Inter")
    # Two big blocks
    add_rounded(s, Inches(0.6), Inches(3.2), Inches(6.0), Inches(3.6), CREAM)
    add_text(s, Inches(0.9), Inches(3.4), Inches(5.4), Inches(0.4),
             "META ADS · 60%", size=11, bold=True, color=BURNT, font="Inter")
    add_text(s, Inches(0.9), Inches(3.8), Inches(5.4), Inches(0.8),
             "€1,200", size=44, bold=True, color=FOREST, font="Georgia")
    meta_lines = [
        "Awareness:    €300  · Carousel ads, broad reach",
        "Lead Gen:     €600  · Conversion-optimised, interest targeting",
        "Retargeting:  €300  · Dynamic ads for site visitors",
        "",
        "Expected: 150–180k impressions · 1.5–2k clicks",
        "3.5–4k LP visits · 200–250 leads at ~€6 CPL",
    ]
    add_text(s, Inches(0.9), Inches(4.7), Inches(5.4), Inches(2.0),
             "\n".join(meta_lines), size=10, color=FOREST, font="Courier New")
    add_rounded(s, Inches(6.7), Inches(3.2), Inches(6.0), Inches(3.6), CREAM)
    add_text(s, Inches(7.0), Inches(3.4), Inches(5.4), Inches(0.4),
             "GOOGLE ADS · 40%", size=11, bold=True, color=BURNT, font="Inter")
    add_text(s, Inches(7.0), Inches(3.8), Inches(5.4), Inches(0.8),
             "€800", size=44, bold=True, color=FOREST, font="Georgia")
    google_lines = [
        "Branded:    €100  · Protect brand-name searches",
        "Generic:    €500  · \"weekend Barcellona\", \"viaggi economici\"",
        "Competitor: €200  · \"volagratis alternative\"",
        "",
        "Expected: 400–700 clicks at €1.50–2.50 CPC",
        "1–1.5k LP visits · 70–100 leads at ~€8 CPL",
    ]
    add_text(s, Inches(7.0), Inches(4.7), Inches(5.4), Inches(2.0),
             "\n".join(google_lines), size=10, color=FOREST, font="Courier New")
    footer(s, idx, total, dark=True)


@slide
def s_meta_detail(s, idx, total):
    header(s, "Meta Ads — campaign structure", "Italian millennials spend 2–3 hours daily on Instagram & Facebook")
    rows = [
        ("Brand Awareness", "€300", "Italian millennials 25–39, urban, travel interests + lookalike of website visitors", "Carousel ads showcasing 5 destinations", "Reach + Engagement"),
        ("Lead Generation", "€600", "Conversion-optimised; interests: travel, weekend trips, food & wine, lookalikes", "Single-image ads + video; CTA: Get the coupon", "Coupon Lead conversion"),
        ("Retargeting", "€300", "Site visitors past 30 days who didn't submit form", "Dynamic ads with destination viewed + €50 reminder", "Coupon Lead conversion"),
    ]
    add_table(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(3.5),
              ["Sub-campaign", "Budget", "Audience", "Creative", "Optimisation"], rows, font_size=10)
    add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
             "Italy Meta CPM €6–9 · CPC ~€1.05 · Travel CPC even lower (€0.63–1.86). Visual platform fits destination marketing.",
             size=11, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_google_detail(s, idx, total):
    header(s, "Google Ads — campaign structure", "Capture active intent · branded protection · competitor displacement")
    rows = [
        ("Branded", "€100", "\"backpack & bottle\", \"backpack and bottle\"", "Defensive — protect from competitor bidding", "Branded conversions"),
        ("Generic Travel", "€500", "\"weekend a Barcellona\", \"viaggi economici Europa\", \"city break Praga\"", "Capture high-intent research", "Lead conversion"),
        ("Competitor", "€200", "\"volagratis alternativa\", \"lastminute recensioni\"", "Win audiences disillusioned with budget OTAs", "Lead conversion"),
    ]
    add_table(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(3.0),
              ["Sub-campaign", "Budget", "Keywords", "Strategy", "Optimisation"], rows, font_size=10)
    add_text(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.4),
             "Expected: 400–700 clicks · €1.50–2.50 CPC · 70–100 leads at ~€8 CPL.",
             size=11, color=FOREST, font="Inter")
    add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.5),
             "Why higher CPL on Google: search intent users are more price-comparing. We accept higher CPL because they convert to bookings at 2× the rate of social leads.",
             size=11, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_channels_excluded(s, idx, total):
    header(s, "Channels we did not use", "And why excluding them is a strategic choice")
    rows = [
        ("Display / programmatic", "Too low intent, high wastage on banner ads"),
        ("TikTok paid", "Audience skews younger (Gen Z) with lower booking propensity"),
        ("Influencer partnerships", "Outside €2,000 budget; would dilute the test"),
        ("Affiliate / referral", "Requires existing brand awareness we don't have yet"),
        ("Print / OOH", "Audience is digital-first; ROI invisible at this scale"),
        ("Email purchase lists", "GDPR violation + brand reputation risk"),
    ]
    add_table(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(3.5),
              ["Channel", "Reason for exclusion"], rows, font_size=12)
    footer(s, idx, total)


@slide
def s_budget(s, idx, total):
    header(s, "Budget breakdown & pacing", dark=True)
    add_rounded(s, Inches(0.6), Inches(2.4), Inches(6.0), Inches(4.5), CREAM)
    add_text(s, Inches(0.9), Inches(2.6), Inches(5.4), Inches(0.4),
             "ALLOCATION", size=10, bold=True, color=BURNT, font="Inter")
    rows = [
        ("Meta Ads", "€1,200", "60%"),
        ("Google Ads", "€800", "40%"),
        ("Total Paid", "€2,000", "100%"),
        ("Owned/Organic", "€0", "Internal time"),
    ]
    y0 = Inches(3.05)
    for i, (a, b, c) in enumerate(rows):
        yy = y0 + Inches(i * 0.45)
        bold = i == 2
        add_text(s, Inches(0.9), yy, Inches(2.5), Inches(0.4), a, size=12, bold=bold, color=FOREST, font="Inter")
        add_text(s, Inches(3.4), yy, Inches(1.8), Inches(0.4), b, size=12, bold=bold, color=FOREST, font="Inter")
        add_text(s, Inches(5.0), yy, Inches(1.5), Inches(0.4), c, size=12, bold=bold, color=BURNT, font="Inter")
    add_text(s, Inches(0.9), Inches(5.0), Inches(5.4), Inches(0.4),
             "ADDITIONAL COSTS", size=10, bold=True, color=BURNT, font="Inter")
    add_text(s, Inches(0.9), Inches(5.4), Inches(5.4), Inches(1.4),
             "Domain + hosting:  €0  (Vercel hobby + .vercel.app)\nEmail platform:    €0–50/mo  (Mailchimp free / Klaviyo)\nStock images:      €0  (Unsplash free licence)\nCreative production: €0  (in-house)",
             size=10, color=FOREST, font="Courier New")
    add_rounded(s, Inches(6.7), Inches(2.4), Inches(6.0), Inches(4.5), CREAM)
    add_text(s, Inches(7.0), Inches(2.6), Inches(5.4), Inches(0.4),
             "PACING (6 WEEKS)", size=10, bold=True, color=BURNT, font="Inter")
    pacing = [
        ("Week 1–2", "€300/wk", "Testing — multiple audiences + creatives"),
        ("Week 3–4", "€400/wk", "Scaling — double down on winners"),
        ("Week 5–6", "€300/wk", "Optimisation — efficiency push"),
    ]
    y = Inches(3.05)
    for w, b, note in pacing:
        add_text(s, Inches(7.0), y, Inches(2.0), Inches(0.4), w, size=12, bold=True, color=FOREST, font="Inter")
        add_text(s, Inches(8.6), y, Inches(1.4), Inches(0.4), b, size=12, color=BURNT, font="Inter")
        add_text(s, Inches(7.0), y + Inches(0.35), Inches(5.0), Inches(0.5), note, size=10, color=FOREST, font="Inter")
        y += Inches(0.95)
    add_text(s, Inches(7.0), Inches(6.0), Inches(5.4), Inches(0.7),
             "ROI framework: 300 leads × 10% nurture-to-booking × €400 avg booking × 10% commission = €1,200 directly attributable revenue. Lead gen — true ROI realised over 90-day coupon validity.",
             size=10, color=FOREST, font="Inter")
    footer(s, idx, total, dark=True)


@slide
def s_timeline(s, idx, total):
    header(s, "Campaign timeline", "8 weeks total · Launch 5 May · Active 5 May – 15 June 2026")
    rows = [
        ("Pre-launch · Week 1 (Apr 1–7)", "Foundation", "Strategy lock, GTM container creation, GA4 property, Meta Business Manager + Pixel install, Google Ads account"),
        ("Pre-launch · Week 2 (Apr 8–14)", "Content", "Landing page copy, hero image selection, blog posts, ad creative, email templates"),
        ("Pre-launch · Week 3 (Apr 15–21)", "Integration", "Full tracking QA, form-to-ESP wire-up, double opt-in, Looker dashboard build"),
        ("Pre-launch · Week 4 (Apr 22–30)", "Testing", "Cross-device QA, dataLayer event verification, conversion test, soft-launch to small audience"),
        ("Active · Week 1 (May 5–11)", "Soft launch", "Monitor closely, gather data, fix issues fast"),
        ("Active · Week 2 (May 12–18)", "Data collection", "Identify winning audiences and creatives"),
        ("Active · Week 3 (May 19–25)", "Scaling", "Pause losers, double budget on winners"),
        ("Active · Week 4 (May 26–Jun 1)", "Peak performance", "Maximise volume at target CPL"),
        ("Active · Week 5 (Jun 2–8)", "Optimisation", "Fine-tune efficiency, lower CPL"),
        ("Active · Week 6 (Jun 9–15)", "Final push", "Hit goals, exhaust budget"),
        ("Post · Week 7 (Jun 16–22)", "Analysis", "Final report, KPI verdict, learnings"),
    ]
    add_table(s, Inches(0.6), Inches(2.4), Inches(12.1), Inches(4.5),
              ["Phase · dates", "Theme", "Activities"], rows, font_size=9)
    footer(s, idx, total)


@slide
def s_email_flows(s, idx, total):
    header(s, "Email automation", "What happens after the user submits the form", dark=True)
    flows = [
        ("T+0 min", "Coupon delivery", "BB50 + PDF guide attached. Confirms double opt-in if required."),
        ("T+1 day", "Welcome email", "Brand intro · 5 destinations briefly · soft CTA to first city."),
        ("T+3 days", "Destination spotlight", "Deep dive on the destination they selected as 'most curious'."),
        ("T+7 days", "Booking psychology", "Why booking 6–8 weeks ahead is the sweet spot (links to blog post)."),
        ("Weekly (Tue 10am)", "Deal alert", "One curated package, one reason to act now, transparent pricing."),
        ("T-30, T-7, T-1 days before expiry", "Coupon expiry reminders", "Soft reminder · stronger urgency · final day."),
        ("T+30 days, no booking", "Re-engagement", "Different angle: \"What's stopping you?\" with FAQ + support contact."),
    ]
    rows = [(f, t, b) for (f, t, b) in flows]
    add_table(s, Inches(0.6), Inches(2.4), Inches(12.1), Inches(4.4),
              ["When", "Email", "Content"], rows,
              header_fill=BURNT, header_text=CREAM, body_text=FOREST, font_size=11)
    add_text(s, Inches(0.6), Inches(6.9), Inches(12), Inches(0.4),
             "Tooling: Mailchimp (free up to 500) for student build; would migrate to Klaviyo at scale for behaviour-triggered flows.",
             size=10, color=CREAM, font="Inter")
    footer(s, idx, total, dark=True)


@slide
def s_dashboard(s, idx, total):
    header(s, "Looker Studio dashboard", "What we look at every morning during the campaign")
    panels = [
        ("Acquisition", "Sessions, users, source/medium, UTMs", "Compare paid vs organic"),
        ("Funnel", "Visit → form_start → form_submit → coupon_download", "Step-by-step drop-off"),
        ("Cost", "Spend, CPC, CPL by channel + sub-campaign", "Spot inefficient spend"),
        ("Audience", "Geo (Italian regions), device, age, gender", "Find concentration"),
        ("Engagement", "Avg session duration, scroll depth, blog reads", "Quality of traffic"),
        ("Coupon", "Downloads by destination · redemption rate", "Which cities drive bookings"),
    ]
    y = Inches(2.5)
    for i, (h, b, c) in enumerate(panels):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.15)
        yy = y + Inches(row * 1.95)
        add_rounded(s, x, yy, Inches(3.95), Inches(1.8), WHITE)
        add_rect(s, x, yy, Inches(3.95), Pt(3), BURNT)
        add_text(s, x + Inches(0.25), yy + Inches(0.2), Inches(3.5), Inches(0.4),
                 h.upper(), size=10, bold=True, color=BURNT, font="Inter")
        add_text(s, x + Inches(0.25), yy + Inches(0.55), Inches(3.5), Inches(0.5),
                 b, size=11, bold=True, color=FOREST, font="Inter")
        add_text(s, x + Inches(0.25), yy + Inches(1.15), Inches(3.5), Inches(0.5),
                 c, size=10, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_assumptions(s, idx, total):
    header(s, "Assumptions made explicit", "Things the brief left open — and what we decided")
    rows = [
        ("Discount amount", "€50 off bookings €300+ — 16.7% meaningful, doesn't devalue brand"),
        ("Validity period", "90 days — matches Italian booking cycle"),
        ("Travel dates", "1 Jun – 31 Dec 2026, excl. Christmas, NYE, Ferragosto"),
        ("Geographic scope", "Italian residents only · landing page in Italian (EN toggle for academic submission)"),
        ("Refund policy", "Standard airline cancellation per carrier"),
        ("Payment", "Card · PayPal · immediate e-ticket delivery"),
        ("Eligibility", "18+, valid email required, phone optional"),
        ("Email frequency", "1 welcome series + 1 weekly deal + expiry reminders — opt-out always one click"),
    ]
    add_table(s, Inches(0.6), Inches(2.4), Inches(12.1), Inches(4.4),
              ["Decision area", "What we assumed"], rows, font_size=11)
    footer(s, idx, total)


@slide
def s_risks(s, idx, total):
    header(s, "Risks & mitigations", dark=True)
    rows = [
        ("CPL exceeds €6.67", "Pause weakest sub-campaign weekly; reallocate to retargeting (lowest CPL)"),
        ("Coupon redemption < 10%", "A/B test discount value (€50 vs €75) in week 4; longer email nurture sequence"),
        ("Italian Q2 holiday clusters", "Timeline avoids Easter, Liberation Day, Labour Day, Republic Day"),
        ("Privacy regulator scrutiny", "Consent Mode v2 default-denied, DPO email visible, retention policy enforced"),
        ("Single-channel dependency", "Two paid channels + organic blog SEO + email — no single point of failure"),
        ("Brand-name confusion", "Branded search budget (€100) protects \"backpack & bottle\" SERPs"),
    ]
    add_table(s, Inches(0.6), Inches(2.4), Inches(12.1), Inches(4.0),
              ["Risk", "Mitigation"], rows,
              header_fill=BURNT, header_text=CREAM, body_text=FOREST, font_size=11)
    footer(s, idx, total, dark=True)


@slide
def s_internal_check(s, idx, total):
    header(s, "Internal consistency check", "Audience → Offer → Positioning → Media → Budget → KPIs → Timing all align")
    chains = [
        ("Audience → Offer", "Millennials value experiences → city break packages with authentic local content"),
        ("Audience → Coupon", "Budget-conscious but not cheap → €50 discount, meaningful but not feels cheap"),
        ("Audience → Curation", "Decision fatigue → 3–5 curated options per city, not 500"),
        ("Offer → Positioning", "Curated packages → \"Skip the research fatigue\" message"),
        ("Offer → Trust", "Transparent pricing → differentiation from sketchy OTAs"),
        ("Positioning → Media", "\"Instagram-worthy\" → visual Meta ads · \"Planning fatigue\" → search ads on intent"),
        ("Media → Budget", "Meta-heavy where audience spends time + affordable Italian CPCs"),
        ("Budget → KPIs", "€6.67 target CPL achievable with Meta €0.70–1.00 CPC + 6% form CVR"),
        ("KPIs → Timing", "May launch hits the peak summer-trip booking window"),
    ]
    y = Inches(2.4)
    for i, (h, b) in enumerate(chains):
        yy = y + Inches(i * 0.48)
        add_text(s, Inches(0.6), yy, Inches(3.4), Inches(0.4), h, size=12, bold=True, color=BURNT, font="Inter")
        add_text(s, Inches(4.1), yy, Inches(8.6), Inches(0.4), b, size=11, color=FOREST, font="Inter")
    add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             "Verdict: no contradictions. Brief is internally coherent and execution-ready.",
             size=12, bold=True, color=BURNT, font="Inter")
    footer(s, idx, total)


@slide
def s_deliverables(s, idx, total):
    header(s, "What we delivered", "Beyond the brief: an actual working campaign", dark=True)
    items = [
        ("Live website", "Next.js + Vercel · /vercel.app subdomain · password-protected"),
        ("Brand identity system", "Wordmark + line mark · 6-colour palette · Fraunces + Inter typography · voice guide"),
        ("Branded coupon PDF", "BB50 single-page PDF generated programmatically — served from /public"),
        ("Public measurement page", "/measurement on the live site documents the entire tracking stack"),
        ("3 SEO blog posts", "Weekend itinerary · booking timing · destination comparison"),
        ("Italian + English landing", "EN/IT toggle · localStorage persisted · academic-friendly"),
        ("GDPR Consent Mode v2", "Default-denied banner · audit-ready · all 4 signals modelled"),
        ("Tracking dataLayer", "10 custom events · 10 GTM tags · 2 custom dimensions · ready to publish"),
        ("This deck", "30 slides covering brief → strategy → brand → site → tracking → media → KPIs"),
    ]
    y = Inches(2.4)
    for i, (h, b) in enumerate(items):
        col = i // 5
        row = i % 5
        x = Inches(0.6 + col * 6.3)
        yy = y + Inches(row * 0.85)
        add_text(s, x, yy, Inches(0.4), Inches(0.4), "✓", size=18, bold=True, color=BURNT, font="Inter")
        add_text(s, x + Inches(0.5), yy, Inches(5.5), Inches(0.4), h, size=12, bold=True, color=CREAM, font="Inter")
        add_text(s, x + Inches(0.5), yy + Inches(0.35), Inches(5.5), Inches(0.5), b, size=10, color=CREAM, font="Inter")
    footer(s, idx, total, dark=True)


@slide
def s_next_steps(s, idx, total):
    header(s, "Next steps to go truly live", "What a real launch would require beyond this academic build")
    steps = [
        ("Legal entity + payment", "VAT registration · Stripe Italia · IATA accreditation or partner OTA contract"),
        ("Real supply", "Negotiate flight + hotel partnership rates for the 5 destinations"),
        ("Booking engine", "Replace static destination cards with a real booking flow + inventory feed"),
        ("Customer support", "Italian-language email + chat with 24h response SLA"),
        ("Email scale", "Klaviyo with behaviour-triggered flows (e.g., abandoned form, viewed-but-no-form)"),
        ("Performance creative", "Ongoing weekly creative refresh — Meta fatigue at 3–5 days for cold audiences"),
        ("Real Looker connection", "GA4 → BigQuery export → Looker Studio for unsampled long-term reporting"),
        ("Compliance audit", "External DPO review pre-launch · cookie audit · privacy policy lawyer review"),
    ]
    y = Inches(2.4)
    for i, (h, b) in enumerate(steps):
        yy = y + Inches(i * 0.55)
        add_text(s, Inches(0.6), yy, Inches(0.6), Inches(0.4), f"0{i+1}", size=14, bold=True, color=BURNT, font="Georgia")
        add_text(s, Inches(1.3), yy, Inches(3.5), Inches(0.4), h, size=12, bold=True, color=FOREST, font="Inter")
        add_text(s, Inches(4.9), yy, Inches(7.8), Inches(0.5), b, size=11, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_links(s, idx, total):
    header(s, "Where to find everything", dark=True)
    items = [
        ("Live site", "https://[deployment].vercel.app  (password-protected — see HANDOFF.md)"),
        ("GitHub repo", "https://github.com/saikanagat1117-sys/backpack-bottle"),
        ("Measurement page (live)", "https://[deployment].vercel.app/measurement"),
        ("Branded coupon PDF", "https://[deployment].vercel.app/backpack-bottle-coupon-BB50.pdf"),
        ("Privacy & T&Cs", "https://[deployment].vercel.app/privacy"),
        ("Blog", "https://[deployment].vercel.app/blog"),
        ("This deck", "Backpack-Bottle-Campaign.pptx (in repo /public folder)"),
    ]
    y = Inches(2.6)
    for i, (h, b) in enumerate(items):
        yy = y + Inches(i * 0.6)
        add_text(s, Inches(0.6), yy, Inches(3.4), Inches(0.4), h, size=12, bold=True, color=BURNT, font="Inter")
        add_text(s, Inches(4.1), yy, Inches(8.6), Inches(0.4), b, size=11, color=CREAM, font="Courier New")
    footer(s, idx, total, dark=True)


@slide
def s_clarity(s, idx, total):
    header(s, "Behavioural insight: Microsoft Clarity", "Heatmaps + session recordings + rage clicks — free, unlimited, GDPR-compliant", dark=True)
    add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             "Where GA4 tells you 'what happened', Clarity tells you 'why'. We added it on top of GA4 — same consent gate, no extra cookies for the user.",
             size=12, color=CREAM, font="Inter")
    rows = [
        ("Click heatmaps", "Where users click on every page", "Find dead zones + ignored CTAs"),
        ("Move heatmaps", "Mouse-tracking attention map", "Where attention concentrates"),
        ("Scroll heatmaps", "How far users actually scroll", "Decide what content lives above the fold"),
        ("Session recordings", "Replay full anonymised sessions", "Diagnose form abandonment + UX bugs"),
        ("Rage clicks", "3+ clicks in same area within 1s", "Frustration moments flagged automatically"),
        ("Dead clicks", "Clicks that didn't change DOM", "Broken interactions — false affordances"),
        ("Quick backs", "Land + leave under 3s", "Bounce-rate quality signal"),
        ("Smart segments", "Auto-built audience cohorts", "Find your power users instantly"),
    ]
    add_table(s, Inches(0.6), Inches(3.1), Inches(12.1), Inches(3.6),
              ["Signal", "What it captures", "What we do with it"], rows,
              header_fill=BURNT, header_text=CREAM, body_text=FOREST, font_size=10)
    add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
             "Plus: we mirror our custom dataLayer events into Clarity tags, so we can filter session recordings by 'rage_click users' or 'high-intent leads who didn't submit.'",
             size=11, color=CREAM, font="Inter")
    footer(s, idx, total, dark=True)


@slide
def s_full_event_matrix(s, idx, total):
    header(s, "Full interaction tracking matrix", "40+ trackable events across 9 categories")
    rows = [
        ("Acquisition", "page_view, visit_count_set, site_search, language_toggle"),
        ("Engagement", "scroll_50, scroll_75, time_on_page (30/60/120/300s), engaged_session"),
        ("Intent / CTAs", "cta_click, nav_click"),
        ("Catalogue", "destination_card_click, destination_detail_view, destination_filter, related_destination_click, gallery_view"),
        ("Decision", "comparison_open, comparison_add, comparison_remove, pricing_filter, currency_change"),
        ("Form / lead", "form_view, form_start, form_field_focus, form_field_blur, form_validation_error, form_submit, form_submit_error, coupon_download, thank_you_view, newsletter_signup"),
        ("Content", "faq_open"),
        ("Sharing / outbound", "share_click, outbound_click, pdf_download, mailto_click, tel_click"),
        ("Behavioural", "rage_click, dead_click, console_error"),
        ("Privacy", "consent_update"),
        ("Media (optional)", "video_play, video_pause, video_complete"),
    ]
    add_table(s, Inches(0.6), Inches(2.4), Inches(12.1), Inches(4.4),
              ["Category", "Events"], rows, font_size=10)
    add_text(s, Inches(0.6), Inches(6.9), Inches(12), Inches(0.4),
             "Live glossary: backpack-bottle.vercel.app/tracking",
             size=11, bold=True, color=BURNT, font="Inter")
    footer(s, idx, total)


@slide
def s_destination_detail(s, idx, total):
    header(s, "Destination detail experience", "10 cities, each a fully realised page")
    add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             "Each destination has its own URL: /destinazioni/[slug]. Cinematic hero, intro, gallery, neighborhoods, day-by-day itinerary, food picks, what's included, weather + budget tips, FAQ, related cities.",
             size=12, color=FOREST, font="Inter")
    cities = ["Roma 🇮🇹", "Barcellona 🇪🇸", "Amsterdam 🇳🇱", "Lisbona 🇵🇹", "Praga 🇨🇿", "Parigi 🇫🇷", "Vienna 🇦🇹", "Berlino 🇩🇪", "Copenaghen 🇩🇰", "Atene 🇬🇷"]
    y = Inches(3.2)
    for i, c in enumerate(cities):
        col = i % 5
        row = i // 5
        x = Inches(0.6 + col * 2.55)
        yy = y + Inches(row * 0.8)
        add_rounded(s, x, yy, Inches(2.4), Inches(0.65), CREAM)
        add_text(s, x + Inches(0.2), yy + Inches(0.18), Inches(2.0), Inches(0.4),
                 c, size=12, bold=True, color=FOREST, font="Inter")
    sections = [
        "Hero with cinematic image + flag + price + flight time + season",
        "Intro: why this city, why our package, neighborhoods picked",
        "Gallery (3 curated images with click tracking)",
        "Day-by-day itinerary with timed suggestions",
        "Three food addresses you can't miss",
        "What's included + budget tip + weather hint",
        "FAQ accordion (city-specific)",
        "Sticky CTA: 'Get €50 off [city]'",
        "Related cities cross-sell",
    ]
    add_text(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
             "EVERY DETAIL PAGE INCLUDES", size=10, bold=True, color=BURNT, font="Inter")
    add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(1.5),
             " • " + "\n • ".join(sections), size=10, color=FOREST, font="Inter")
    footer(s, idx, total)


@slide
def s_zapier_flow(s, idx, total):
    header(s, "Lead-to-customer automation chain", "Zapier orchestrates the moment after submit", dark=True)
    add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             "Form submit → Vercel API → Zapier webhook → 5-way fan-out: Mailchimp, Sheets, Slack, Meta CAPI, internal log.",
             size=12, color=CREAM, font="Inter")
    diagram = (
        "                              POST /api/lead\n"
        "                                    │\n"
        "                                    ▼\n"
        "                          Zapier 'Catch B&B lead'\n"
        "                                    │\n"
        "          ┌─────────┬─────────┬─────────┬──────────┐\n"
        "          ▼           ▼           ▼           ▼            ▼\n"
        "      Mailchimp     Google      Slack       Meta CAPI    Internal\n"
        "      (welcome      Sheets      #bb-leads   server-side  Vercel\n"
        "       series)      master      alert       Lead event   logs\n"
        "          │\n"
        "          ▼\n"
        "    T+0:  Coupon delivery\n"
        "    T+1d: Welcome + brand intro\n"
        "    T+3d: Destination deep-dive (personalised on 'interest')\n"
        "    T+7d: Booking psychology\n"
        "    Tue:  Weekly deal alert\n"
        "    T+30: Reminder if no booking\n"
        "    T+83: Expiry warning (7 days left)\n"
        "    T+45: Re-engage if dormant"
    )
    add_rounded(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(4.0), FOREST_DK)
    add_text(s, Inches(0.9), Inches(3.1), Inches(11.5), Inches(3.9),
             diagram, size=9, color=CREAM, font="Courier New")
    footer(s, idx, total, dark=True)


@slide
def s_thanks(s, idx, total):
    add_bg(s, FOREST)
    add_rect(s, Inches(0.6), Inches(0.55), Inches(0.4), Pt(2), BURNT)
    add_text(s, Inches(1.1), Inches(0.45), Inches(11), Inches(0.4),
             "GROUP 3 · DIGITAL PLATFORMS LAB", size=10, bold=True, color=BURNT, font="Inter")
    add_text(s, Inches(0.6), Inches(2.5), Inches(12), Inches(2.0),
             "Thank you.", size=80, bold=True, color=CREAM, font="Georgia")
    add_text(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.6),
             "Skip the research fatigue.", size=22, color=BURNT, font="Georgia")
    add_text(s, Inches(0.6), Inches(6.8), Inches(8), Inches(0.4),
             "Sai Prathyaksh Kanagat · Group 3 · Bologna Business School",
             size=12, color=CREAM, font="Inter")
    add_text(s, Inches(8.5), Inches(6.8), Inches(4.4), Inches(0.4),
             "saikanagat1117@gmail.com", size=12, color=CREAM, font="Inter", align=PP_ALIGN.RIGHT)


# ---- build ----
total = len(slides_data)
for i, fn in enumerate(slides_data, start=1):
    s = prs.slides.add_slide(BLANK)
    fn(s, i, total)

out = "Backpack-Bottle-Campaign.pptx"
prs.save(out)
print(f"Saved {out} with {total} slides")
